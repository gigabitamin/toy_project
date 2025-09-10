#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Puppeteer HTML to PPTX Converter
Puppeteer를 사용하여 페이지 완전 로딩 후 스크린샷을 찍는 변환기
"""

import os
import sys
import subprocess
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import tempfile
import shutil
from bs4 import BeautifulSoup

class PuppeteerHTMLToPPTXConverter:
    def __init__(self, html_dir, output_path):
        self.html_dir = Path(html_dir)
        self.output_path = Path(output_path)
        self.temp_dir = None
    
    def calculate_content_height(self, html_content):
        """HTML 내용의 실제 높이를 추정"""
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # 텍스트 내용 길이 기반으로 높이 추정
        text_content = soup.get_text()
        text_length = len(text_content)
        
        # 이미지 개수 확인
        images = soup.find_all('img')
        image_count = len(images)
        
        # 테이블 개수 확인
        tables = soup.find_all('table')
        table_count = len(tables)
        
        # 기본 높이 계산 (텍스트 길이 기반)
        base_height = 800  # 기본 높이
        text_height = text_length * 0.5  # 텍스트당 0.5px
        image_height = image_count * 200  # 이미지당 200px
        table_height = table_count * 150  # 테이블당 150px
        
        estimated_height = int(base_height + text_height + image_height + table_height)
        
        # 최소/최대 높이 제한
        min_height = 1080
        max_height = 2160
        
        final_height = max(min_height, min(estimated_height, max_height))
        
        print(f"내용 분석: 텍스트 {text_length}자, 이미지 {image_count}개, 테이블 {table_count}개")
        print(f"추정 높이: {final_height}px")
        
        return final_height
    
    def adjust_html_height(self, html_content, target_height_px=None):
        """HTML 내용의 높이를 자동으로 조절하여 잘리지 않도록 함"""
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # 목표 높이가 지정되지 않은 경우 자동 계산
        if target_height_px is None:
            target_height_px = self.calculate_content_height(html_content)
        
        # body 태그에 최소 높이 설정
        body = soup.find('body')
        if body:
            # 기존 스타일 확인
            existing_style = body.get('style', '')
            
            # 최소 높이 설정 (픽셀 단위)
            min_height_style = f"min-height: {target_height_px}px; height: auto;"
            
            if existing_style:
                # 기존 스타일과 병합
                body['style'] = f"{existing_style}; {min_height_style}"
            else:
                body['style'] = min_height_style
        
        # html 태그에도 높이 설정
        html_tag = soup.find('html')
        if html_tag:
            existing_style = html_tag.get('style', '')
            html_style = f"height: auto; min-height: {target_height_px}px;"
            
            if existing_style:
                html_tag['style'] = f"{existing_style}; {html_style}"
            else:
                html_tag['style'] = html_style
        
        # CSS 스타일 추가
        style_tag = soup.find('style')
        if not style_tag:
            style_tag = soup.new_tag('style')
            soup.head.append(style_tag)
        
        # 기존 CSS에 높이 관련 스타일 추가
        additional_css = f"""
        /* 자동 높이 조절을 위한 추가 스타일 */
        html, body {{
            height: auto !important;
            min-height: {target_height_px}px !important;
            overflow-x: hidden;
            overflow-y: visible;
            margin: 0;
            padding: 0;
        }}
        
        .page-body, .page {{
            height: auto !important;
            min-height: {target_height_px}px !important;
            margin: 0;
            padding: 20px;
            box-sizing: border-box;
        }}
        
        /* 모든 컨테이너 요소의 높이 자동 조절 */
        div, section, article, main {{
            height: auto !important;
            min-height: fit-content;
        }}
        
        /* 테이블 높이 자동 조절 */
        table {{
            height: auto !important;
            min-height: fit-content;
        }}
        
        /* 이미지 높이 자동 조절 */
        img {{
            max-width: 100%;
            height: auto;
        }}
        """
        
        if style_tag.string:
            style_tag.string += additional_css
        else:
            style_tag.string = additional_css
        
        return str(soup)
    
    def setup_temp_directory(self):
        """임시 디렉토리 설정"""
        self.temp_dir = Path(tempfile.mkdtemp())
        print(f"임시 디렉토리 생성: {self.temp_dir}")
        
    def cleanup_temp_directory(self):
        """임시 디렉토리 정리"""
        if self.temp_dir and self.temp_dir.exists():
            shutil.rmtree(self.temp_dir)
            print("임시 디렉토리 정리 완료")
    
    def convert_html_to_image_puppeteer(self, html_file, slide_number):
        """Puppeteer를 사용하여 HTML 파일을 이미지로 변환"""
        try:
            # HTML 파일 경로
            html_path = self.html_dir / html_file
            
            # 출력 이미지 경로
            output_image = self.temp_dir / f"slide_{slide_number:02d}.png"
            
            print(f"변환 중: {html_file} -> {output_image.name}")
            
            # HTML 파일 읽기
            with open(html_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # HTML 높이 자동 조절
            adjusted_html = self.adjust_html_height(html_content)
            
            # 조절된 높이 다시 계산
            calculated_height = self.calculate_content_height(adjusted_html)
            print(f"최종 이미지 높이: {calculated_height}px")
            
            # 임시 HTML 파일 생성
            temp_html_path = self.temp_dir / f"temp_{html_file}"
            with open(temp_html_path, 'w', encoding='utf-8') as f:
                f.write(adjusted_html)
            
            # Puppeteer 스크립트 생성 (절대 경로 사용)
            puppeteer_script = f"""
const puppeteer = require('C:/Project/gigabitamin/genspark/node_modules/puppeteer');
const path = require('path');

(async () => {{
    const browser = await puppeteer.launch({{ 
        headless: true,
        args: ['--no-sandbox', '--disable-setuid-sandbox']
    }});
    const page = await browser.newPage();
    
    // 뷰포트 설정
    await page.setViewport({{ width: 1920, height: {calculated_height} }});
    
    // HTML 파일 로드
    const htmlPath = path.resolve('{temp_html_path.as_posix()}');
    await page.goto(`file://${{htmlPath}}`, {{ 
        waitUntil: 'networkidle0'  // 모든 네트워크 요청 완료까지 대기
    }});
    
    // 추가 대기 (폰트와 CSS 완전 로딩)
    await new Promise(resolve => setTimeout(resolve, 3000));
    
    // 스크린샷 촬영
    await page.screenshot({{
        path: '{output_image.as_posix()}',
        fullPage: true,
        type: 'png'
    }});
    
    await browser.close();
    console.log('스크린샷 완료: {output_image.name}');
}})();
"""
            
            # Puppeteer 스크립트 파일 생성
            script_path = self.temp_dir / f"screenshot_{slide_number}.js"
            with open(script_path, 'w', encoding='utf-8') as f:
                f.write(puppeteer_script)
            
            print(f"  Puppeteer로 페이지 로딩 대기 중...")
            
            # Puppeteer 실행 (node_modules 경로 지정)
            result = subprocess.run([
                'node', str(script_path)
            ], capture_output=True, text=True, cwd=str(Path('C:/Project/gigabitamin/genspark')))
            
            if result.returncode != 0:
                print(f"Puppeteer 실행 오류: {result.stderr}")
                return None
            
            # 이미지 파일 확인
            if output_image.exists():
                # 이미지 크기 확인
                with Image.open(output_image) as img:
                    actual_width, actual_height = img.size
                    print(f"생성된 이미지 크기: {actual_width}x{actual_height}")
                
                # 이미지 크기 조정 (PPT 슬라이드에 맞게 리사이즈)
                print(f"리사이즈 전 이미지 크기: {actual_width}x{actual_height}")
                self.resize_image_to_fit(output_image)
                
                # 리사이즈 후 크기 확인
                with Image.open(output_image) as resized_img:
                    final_width, final_height = resized_img.size
                    print(f"리사이즈 후 이미지 크기: {final_width}x{final_height}")
                
                return output_image
            else:
                print(f"이미지 파일이 생성되지 않았습니다: {output_image}")
                return None
            
        except Exception as e:
            print(f"HTML 변환 오류 ({html_file}): {e}")
            return None
    
    def resize_image_to_fit(self, image_path):
        """이미지를 1920x1080에 맞게 리사이즈 (PPT 슬라이드 크기)"""
        try:
            from PIL import Image
            
            with Image.open(image_path) as img:
                # 원본 이미지 크기
                original_width, original_height = img.size
                
                print(f"원본 이미지 크기: {original_width}x{original_height}")
                
                # 목표 크기 (PPT 슬라이드 크기)
                target_width = 1920
                target_height = 1080
                
                # 원본 비율 계산
                original_ratio = original_width / original_height
                target_ratio = target_width / target_height
                
                if original_ratio > target_ratio:
                    # 원본이 더 넓은 경우 - 높이를 기준으로 리사이즈
                    new_height = target_height
                    new_width = int(target_height * original_ratio)
                else:
                    # 원본이 더 높은 경우 - 너비를 기준으로 리사이즈
                    new_width = target_width
                    new_height = int(target_width / original_ratio)
                
                print(f"리사이즈 크기: {new_width}x{new_height}")
                
                # 리사이즈
                resized_img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
                
                # 1920x1080으로 크롭 (상단부터 시작)
                if new_width > target_width:
                    # 너비가 더 긴 경우 - 중앙에서 크롭
                    left = (new_width - target_width) // 2
                    right = left + target_width
                else:
                    left = 0
                    right = target_width
                
                if new_height > target_height:
                    # 높이가 더 긴 경우 - 상단부터 크롭 (위쪽이 잘리지 않도록)
                    top = 0
                    bottom = target_height
                else:
                    top = 0
                    bottom = new_height
                
                print(f"크롭 영역: left={left}, top={top}, right={right}, bottom={bottom}")
                
                cropped_img = resized_img.crop((left, top, right, bottom))
                
                # 최종 크기가 1920x1080이 되도록 패딩 추가 (필요시)
                if cropped_img.size != (target_width, target_height):
                    # 흰색 배경으로 새 이미지 생성
                    final_img = Image.new('RGB', (target_width, target_height), 'white')
                    # 중앙에 이미지 붙이기
                    paste_x = (target_width - cropped_img.width) // 2
                    paste_y = (target_height - cropped_img.height) // 2
                    final_img.paste(cropped_img, (paste_x, paste_y))
                    cropped_img = final_img
                
                # 저장
                cropped_img.save(image_path, 'PNG', quality=95)
                
                print(f"이미지 리사이즈 완료: {image_path.name} -> {cropped_img.size}")
                
        except Exception as e:
            print(f"이미지 리사이즈 오류: {e}")
    
    def create_pptx(self, image_files):
        """이미지 파일들로부터 PPTX 생성"""
        try:
            # 새 프레젠테이션 생성
            prs = Presentation()
            
            # 슬라이드 크기 설정 (16:9 비율) - 더 큰 크기로 설정
            prs.slide_width = Inches(20)   # 1920px (1인치 = 75px 기준)
            prs.slide_height = Inches(11.25)  # 1080px
            
            print(f"PPT 슬라이드 크기: {prs.slide_width} x {prs.slide_height}")
            
            for i, image_file in enumerate(image_files):
                if image_file and image_file.exists():
                    # 빈 슬라이드 추가
                    slide_layout = prs.slide_layouts[6]  # 빈 레이아웃
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # 이미지 추가 - 슬라이드 전체 크기로 설정
                    left = Inches(0)
                    top = Inches(0)
                    width = Inches(20)   # 1920px
                    height = Inches(11.25)  # 1080px
                    
                    # 이미지 삽입
                    picture = slide.shapes.add_picture(str(image_file), left, top, width, height)
                    
                    # 이미지가 슬라이드를 완전히 채우도록 설정
                    picture.left = 0
                    picture.top = 0
                    picture.width = prs.slide_width
                    picture.height = prs.slide_height
                    
                    print(f"슬라이드 {i+1} 추가 완료 - 크기: {picture.width} x {picture.height}")
                else:
                    print(f"이미지 파일 없음: {image_file}")
            
            # PPTX 파일 저장
            prs.save(self.output_path)
            print(f"PPTX 파일 저장 완료: {self.output_path}")
            
        except Exception as e:
            print(f"PPTX 생성 오류: {e}")
            raise
    
    def convert(self):
        """전체 변환 프로세스 실행"""
        try:
            # HTML 파일 목록 가져오기
            html_files = sorted([f for f in self.html_dir.glob("*.html")])
            
            if not html_files:
                print("HTML 파일을 찾을 수 없습니다.")
                return False
            
            print(f"발견된 HTML 파일: {len(html_files)}개")
            for html_file in html_files:
                print(f"  - {html_file.name}")
            
            # 임시 디렉토리 설정
            self.setup_temp_directory()
            
            # HTML 파일들을 이미지로 변환
            image_files = []
            for i, html_file in enumerate(html_files, 1):
                image_file = self.convert_html_to_image_puppeteer(html_file.name, i)
                image_files.append(image_file)
            
            # PPTX 생성
            self.create_pptx(image_files)
            
            return True
            
        except Exception as e:
            print(f"변환 프로세스 오류: {e}")
            return False
        
        finally:
            # 임시 디렉토리 정리
            self.cleanup_temp_directory()

def main():
    # 설정
    html_dir = r"C:\Project\gigabitamin\genspark\smart_gate\html"
    output_path = r"C:\Project\gigabitamin\genspark\smart_gate\smart_gate_puppeteer_pptx.pptx"
    
    print("Puppeteer HTML to PPTX 변환기 시작")
    print(f"HTML 디렉토리: {html_dir}")
    print(f"출력 파일: {output_path}")
    print("-" * 50)
    
    # 변환기 생성 및 실행
    converter = PuppeteerHTMLToPPTXConverter(html_dir, output_path)
    success = converter.convert()
    
    if success:
        print("-" * 50)
        print("변환 완료!")
        print(f"출력 파일: {output_path}")
    else:
        print("-" * 50)
        print("변환 실패!")
        sys.exit(1)

if __name__ == "__main__":
    main()
