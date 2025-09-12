#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
CSS-Aware HTML to PPTX Converter
CSS 스타일을 분석하여 디자인 구조를 유지하는 변환기
"""

import os
import sys
import requests
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup
import re
import tempfile
import shutil
from html2image import Html2Image
from PIL import Image

class CSSAwareConverter:
    def __init__(self, html_file, output_path):
        self.html_file = html_file
        self.output_path = output_path
        self.temp_dir = None
        self.hti = Html2Image()
        
    def setup_temp_directory(self):
        """임시 디렉토리 생성"""
        self.temp_dir = Path(tempfile.mkdtemp())
        
    def cleanup_temp_directory(self):
        """임시 디렉토리 정리"""
        if self.temp_dir and self.temp_dir.exists():
            shutil.rmtree(self.temp_dir)
    
    def download_fontawesome_svg(self, icon_class, color='#2563eb'):
        """FontAwesome 아이콘 SVG 다운로드"""
        try:
            clean_class = icon_class.replace('fas ', '').replace('fab ', '').replace('far ', '').replace('fa-', '')
            
            paths = [
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/6.x/svgs/solid/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/6.x/svgs/brands/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/6.x/svgs/regular/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/5.x/svgs/solid/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/5.x/svgs/brands/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/5.x/svgs/regular/{clean_class}.svg",
            ]
            
            for path in paths:
                try:
                    response = requests.get(path, timeout=10)
                    if response.status_code == 200:
                        svg_content = response.text
                        svg_content = svg_content.replace('fill="currentColor"', f'fill="{color}"')
                        svg_content = svg_content.replace('fill="#000"', f'fill="{color}"')
                        svg_content = svg_content.replace('fill="black"', f'fill="{color}"')
                        
                        svg_file = self.temp_dir / f"{clean_class}.svg"
                        with open(svg_file, 'w', encoding='utf-8') as f:
                            f.write(svg_content)
                        return svg_file
                except Exception as e:
                    continue
            return None
        except Exception as e:
            return None
    
    def svg_to_png_with_html2image(self, svg_file, size=64):
        """SVG를 HTML2Image로 PNG 변환"""
        try:
            with open(svg_file, 'r', encoding='utf-8') as f:
                svg_content = f.read()
            
            html_content = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <style>
                    body {{
                        margin: 0;
                        padding: 0;
                        width: {size}px;
                        height: {size}px;
                        display: flex;
                        align-items: center;
                        justify-content: center;
                        background: transparent;
                    }}
                    svg {{
                        width: {size}px;
                        height: {size}px;
                    }}
                </style>
            </head>
            <body>
                {svg_content}
            </body>
            </html>
            """
            
            temp_html = self.temp_dir / f"{svg_file.stem}_temp.html"
            with open(temp_html, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            png_file = self.temp_dir / f"{svg_file.stem}.png"
            self.hti.screenshot(
                html_file=str(temp_html),
                save_as=f"{svg_file.stem}.png",
                size=(size, size)
            )
            
            generated_png = Path(f"{svg_file.stem}.png")
            if generated_png.exists():
                shutil.move(str(generated_png), str(png_file))
                return png_file
            else:
                return None
        except Exception as e:
            return None
    
    def create_text_box(self, slide, text, x, y, width, height, font_size=16, color='#000000', bold=False, align='left', font_family='맑은 고딕'):
        """텍스트 박스 생성"""
        try:
            textbox = slide.shapes.add_textbox(
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            
            text_frame = textbox.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            p = text_frame.paragraphs[0]
            p.text = text
            
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.name = font_family
            
            if color.startswith('#'):
                hex_color = color[1:]
                if len(hex_color) == 3:
                    hex_color = ''.join([c*2 for c in hex_color])
                p.font.color.rgb = RGBColor(
                    int(hex_color[0:2], 16),
                    int(hex_color[2:4], 16),
                    int(hex_color[4:6], 16)
                )
            
            if align == 'center':
                p.alignment = PP_ALIGN.CENTER
            elif align == 'right':
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT
            
            return textbox
        except Exception as e:
            print(f"텍스트 박스 생성 오류: {e}")
            return None
    
    def create_icon_circle(self, slide, icon_class, x, y, size=0.8, bg_color='#dbeafe', icon_color='#2563eb'):
        """아이콘 원형 배경 생성"""
        try:
            # 원형 배경
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x), Inches(y), Inches(size), Inches(size)
            )
            
            # 배경 색상
            fill = circle.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(
                int(bg_color[1:3], 16),
                int(bg_color[3:5], 16),
                int(bg_color[5:7], 16)
            )
            
            # 아이콘 이미지 추가
            if icon_class:
                svg_file = self.download_fontawesome_svg(icon_class, icon_color)
                if svg_file:
                    png_file = self.svg_to_png_with_html2image(svg_file, 32)
                    if png_file and png_file.exists():
                        icon_left = Inches(x + size/4)
                        icon_top = Inches(y + size/4)
                        icon_width = Inches(size/2)
                        icon_height = Inches(size/2)
                        
                        slide.shapes.add_picture(str(png_file), icon_left, icon_top, icon_width, icon_height)
            
            return circle
        except Exception as e:
            print(f"아이콘 원형 생성 오류: {e}")
            return None
    
    def create_tech_card(self, slide, title, subtitle, x, y, width, height, icon_class=None):
        """기술 스택 카드 생성"""
        try:
            # 카드 배경
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            
            # 카드 스타일
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            line = card.line
            line.color.rgb = RGBColor(229, 231, 235)
            line.width = Pt(1)
            
            # 텍스트 프레임
            text_frame = card.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)
            text_frame.margin_top = Inches(0.2)
            text_frame.margin_bottom = Inches(0.2)
            
            # 제목
            p1 = text_frame.paragraphs[0]
            p1.text = title
            p1.alignment = PP_ALIGN.LEFT
            font1 = p1.font
            font1.size = Pt(16)
            font1.bold = True
            font1.color.rgb = RGBColor(31, 41, 55)
            
            # 부제목
            if subtitle:
                p2 = text_frame.add_paragraph()
                p2.text = subtitle
                p2.alignment = PP_ALIGN.LEFT
                font2 = p2.font
                font2.size = Pt(12)
                font2.color.rgb = RGBColor(107, 114, 128)
            
            # 아이콘 추가
            if icon_class:
                self.create_icon_circle(slide, icon_class, x + 0.1, y + 0.1, 0.3, '#f3f4f6', '#6b7280')
            
            return card
        except Exception as e:
            print(f"기술 카드 생성 오류: {e}")
            return None
    
    def create_feature_card(self, slide, title, content, x, y, width, height, icon_class=None):
        """기능 카드 생성"""
        try:
            # 카드 배경
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            
            # 카드 스타일
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(249, 250, 251)
            
            line = card.line
            line.color.rgb = RGBColor(229, 231, 235)
            line.width = Pt(1)
            
            # 텍스트 프레임
            text_frame = card.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)
            text_frame.margin_top = Inches(0.2)
            text_frame.margin_bottom = Inches(0.2)
            
            # 제목
            p1 = text_frame.paragraphs[0]
            p1.text = title
            p1.alignment = PP_ALIGN.LEFT
            font1 = p1.font
            font1.size = Pt(14)
            font1.bold = True
            font1.color.rgb = RGBColor(31, 41, 55)
            
            # 내용
            if content:
                p2 = text_frame.add_paragraph()
                p2.text = content
                p2.alignment = PP_ALIGN.LEFT
                font2 = p2.font
                font2.size = Pt(12)
                font2.color.rgb = RGBColor(107, 114, 128)
            
            # 아이콘 추가
            if icon_class:
                self.create_icon_circle(slide, icon_class, x + 0.1, y + 0.1, 0.2, '#dbeafe', '#2563eb')
            
            return card
        except Exception as e:
            print(f"기능 카드 생성 오류: {e}")
            return None
    
    def create_divider_line(self, slide, x, y, width, color='#3b82f6'):
        """구분선 생성"""
        try:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y), Inches(width), Inches(0.05)
            )
            
            fill = line.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(
                int(color[1:3], 16),
                int(color[3:5], 16),
                int(color[5:7], 16)
            )
            
            return line
        except Exception as e:
            print(f"구분선 생성 오류: {e}")
            return None
    
    def extract_icon_class(self, element):
        """요소에서 아이콘 클래스 추출"""
        if not element:
            return None
            
        icon_elem = element.find('i')
        if icon_elem:
            classes = icon_elem.get('class', [])
            for cls in classes:
                if cls.startswith('fa-'):
                    return cls
        
        classes = element.get('class', [])
        for cls in classes:
            if cls.startswith('fa-'):
                return cls
        
        return None
    
    def parse_html_with_css(self, soup, slide):
        """CSS 스타일을 고려한 HTML 파싱"""
        try:
            y_pos = 0.5
            
            # 01.html 스타일 파싱
            if '01.html' in str(self.html_file):
                self.parse_01_html(soup, slide)
            # 02.html 스타일 파싱
            elif '02.html' in str(self.html_file):
                self.parse_02_html(soup, slide)
            # 03.html 스타일 파싱
            elif '03.html' in str(self.html_file):
                self.parse_03_html(soup, slide)
            # 기타 파일들
            else:
                self.parse_generic_html(soup, slide)
                
        except Exception as e:
            print(f"HTML 파싱 오류: {e}")
    
    def parse_01_html(self, soup, slide):
        """01.html 전용 파싱 (메인 페이지)"""
        try:
            # 메인 제목
            title_elem = soup.find('h1', class_='title')
            if title_elem:
                self.create_text_box(
                    slide, title_elem.get_text().strip(),
                    1.0, 1.0, 11.0, 1.5,
                    font_size=48, color='#2563eb', bold=True, align='center'
                )
            
            # 부제목
            subtitle_elem = soup.find('h2', class_='subtitle')
            if subtitle_elem:
                self.create_text_box(
                    slide, subtitle_elem.get_text().strip(),
                    1.0, 2.8, 11.0, 1.0,
                    font_size=32, color='#1e40af', bold=True, align='center'
                )
            
            # 기술 스택 섹션
            tech_section = soup.find('div', class_='tech-stack-section')
            if tech_section:
                y_pos = 4.0
                
                # 기술 스택 제목
                self.create_text_box(
                    slide, "주요 기술 스택",
                    1.0, y_pos, 11.0, 0.8,
                    font_size=28, color='#1f2937', bold=True, align='center'
                )
                y_pos += 1.0
                
                # 기술 스택 카드들
                tech_cards = tech_section.find_all('div', class_='tech-card')
                for i, card in enumerate(tech_cards[:4]):  # 최대 4개
                    title_elem = card.find('h3')
                    subtitle_elem = card.find('p')
                    icon_elem = card.find('i')
                    
                    title = title_elem.get_text().strip() if title_elem else ""
                    subtitle = subtitle_elem.get_text().strip() if subtitle_elem else ""
                    icon_class = self.extract_icon_class(card)
                    
                    x_pos = 1.0 + (i % 2) * 5.5
                    y_card = y_pos + (i // 2) * 1.5
                    
                    self.create_tech_card(
                        slide, title, subtitle,
                        x_pos, y_card, 5.0, 1.2, icon_class
                    )
            
            # 링크 버튼들
            link_section = soup.find('div', class_='link-section')
            if link_section:
                buttons = link_section.find_all('a', class_='link-button')
                for i, button in enumerate(buttons[:2]):
                    text = button.get_text().strip()
                    x_pos = 2.0 + i * 4.0
                    y_pos = 6.5
                    
                    # 버튼 배경
                    btn = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(x_pos), Inches(y_pos), Inches(3.5), Inches(0.6)
                    )
                    
                    fill = btn.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(37, 99, 235)
                    
                    # 버튼 텍스트
                    self.create_text_box(
                        slide, text,
                        x_pos + 0.1, y_pos + 0.1, 3.3, 0.4,
                        font_size=16, color='#ffffff', bold=True, align='center'
                    )
                    
        except Exception as e:
            print(f"01.html 파싱 오류: {e}")
    
    def parse_02_html(self, soup, slide):
        """02.html 전용 파싱 (프로젝트 개요)"""
        try:
            # 섹션 제목
            title_elem = soup.find('h1', class_='section-title')
            if title_elem:
                self.create_text_box(
                    slide, title_elem.get_text().strip(),
                    1.0, 0.5, 11.0, 1.0,
                    font_size=36, color='#2563eb', bold=True, align='left'
                )
            
            # 구분선
            self.create_divider_line(slide, 1.0, 1.8, 2.0)
            
            # 기능 카드들
            feature_cards = soup.find_all('div', class_='feature-card')
            y_pos = 2.2
            
            for i, card in enumerate(feature_cards[:3]):  # 최대 3개
                title_elem = card.find('h3')
                content_elem = card.find('p')
                icon_elem = card.find('i')
                
                title = title_elem.get_text().strip() if title_elem else ""
                content = content_elem.get_text().strip() if content_elem else ""
                icon_class = self.extract_icon_class(card)
                
                x_pos = 0.5 + (i % 2) * 6.0
                y_card = y_pos + (i // 2) * 2.0
                
                self.create_feature_card(
                    slide, title, content,
                    x_pos, y_card, 5.5, 1.8, icon_class
                )
                
        except Exception as e:
            print(f"02.html 파싱 오류: {e}")
    
    def parse_03_html(self, soup, slide):
        """03.html 전용 파싱 (기술 스택)"""
        try:
            # 메인 제목
            title_elem = soup.find('h1')
            if title_elem:
                self.create_text_box(
                    slide, title_elem.get_text().strip(),
                    1.0, 0.5, 11.0, 1.0,
                    font_size=36, color='#1f2937', bold=True, align='center'
                )
            
            # 기술 스택 카드들
            tech_cards = soup.find_all('div', class_='tech-card')
            y_pos = 2.0
            
            for i, card in enumerate(tech_cards[:4]):  # 최대 4개
                title_elem = card.find('h3')
                subtitle_elem = card.find('p')
                icon_elem = card.find('i')
                
                title = title_elem.get_text().strip() if title_elem else ""
                subtitle = subtitle_elem.get_text().strip() if subtitle_elem else ""
                icon_class = self.extract_icon_class(card)
                
                x_pos = 0.5 + (i % 2) * 6.0
                y_card = y_pos + (i // 2) * 2.0
                
                self.create_tech_card(
                    slide, title, subtitle,
                    x_pos, y_card, 5.5, 1.5, icon_class
                )
                
        except Exception as e:
            print(f"03.html 파싱 오류: {e}")
    
    def parse_generic_html(self, soup, slide):
        """일반적인 HTML 파싱"""
        try:
            y_pos = 0.5
            
            # 제목 찾기
            title_element = soup.find('h1') or soup.find('h2') or soup.find('title')
            if title_element:
                title_text = title_element.get_text().strip()
                if title_text:
                    self.create_text_box(
                        slide, title_text,
                        1.0, y_pos, 11.0, 1.0,
                        font_size=32, color='#1f2937', bold=True, align='center'
                    )
                    y_pos += 1.2
            
            # 텍스트 요소들
            text_elements = soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'div'])
            
            for element in text_elements:
                if y_pos > 6.0:
                    break
                    
                text = element.get_text().strip()
                if not text or len(text) < 3:
                    continue
                
                if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                    icon_class = self.extract_icon_class(element)
                    
                    if icon_class:
                        self.create_icon_circle(slide, icon_class, 0.5, y_pos, 0.4, '#dbeafe', '#2563eb')
                        title_x = 1.0
                    else:
                        title_x = 0.5
                    
                    self.create_text_box(
                        slide, text,
                        title_x, y_pos, 10.0, 0.6,
                        font_size=24, color='#1f2937', bold=True, align='left'
                    )
                    y_pos += 0.8
                
                elif element.name == 'p' and len(text) > 10:
                    self.create_text_box(
                        slide, text,
                        0.5, y_pos, 11.0, 0.8,
                        font_size=16, color='#374151', bold=False, align='left'
                    )
                    y_pos += 1.0
                
        except Exception as e:
            print(f"일반 HTML 파싱 오류: {e}")
    
    def convert(self):
        """HTML을 PPTX로 변환"""
        try:
            self.setup_temp_directory()
            
            with open(self.html_file, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            prs = Presentation()
            prs.slide_width = Inches(13.33)  # 16:9 비율
            prs.slide_height = Inches(7.5)
            
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            
            self.parse_html_with_css(soup, slide)
            
            prs.save(self.output_path)
            print(f"✅ 변환 완료: {Path(self.html_file).name}")
            
            return True
            
        except Exception as e:
            print(f"❌ 변환 실패: {e}")
            return False
        
        finally:
            self.cleanup_temp_directory()

def convert_folder_to_pptx(html_folder, output_path):
    """폴더 내 모든 HTML 파일을 하나의 PPTX로 변환"""
    try:
        html_folder = Path(html_folder)
        if not html_folder.exists():
            print(f"HTML 폴더가 존재하지 않습니다: {html_folder}")
            return False
        
        html_files = list(html_folder.glob("*.html"))
        if not html_files:
            print(f"HTML 파일이 없습니다: {html_folder}")
            return False
        
        print(f"발견된 HTML 파일 {len(html_files)}개:")
        for html_file in html_files:
            print(f"  - {html_file.name}")
        
        prs = Presentation()
        prs.slide_width = Inches(13.33)  # 16:9 비율
        prs.slide_height = Inches(7.5)
        
        for i, html_file in enumerate(html_files):
            print(f"\n--- {html_file.name} 변환 중 ({i+1}/{len(html_files)}) ---")
            
            try:
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                
                converter = CSSAwareConverter(str(html_file), "")
                converter.setup_temp_directory()
                
                with open(html_file, 'r', encoding='utf-8') as f:
                    html_content = f.read()
                
                soup = BeautifulSoup(html_content, 'html.parser')
                converter.parse_html_with_css(soup, slide)
                converter.cleanup_temp_directory()
                
                print(f"✅ {html_file.name} 변환 완료")
                
            except Exception as e:
                print(f"❌ {html_file.name} 변환 실패: {e}")
                continue
        
        prs.save(output_path)
        print(f"\n✅ 모든 HTML 파일이 하나의 PPTX로 변환 완료!")
        print(f"출력 파일: {output_path}")
        return True
        
    except Exception as e:
        print(f"폴더 변환 오류: {e}")
        return False

def main():
    html_folder = r"C:\Project\gigabitamin\genspark\dcs_site\html"
    output_path = r"C:\Project\gigabitamin\genspark\dcs_site\html\css_aware_all_pages.pptx"
    
    print("CSS-Aware HTML to Editable PPTX 변환기")
    print("=" * 50)
    print(f"HTML 폴더: {html_folder}")
    print(f"출력 파일: {output_path}")
    print("-" * 50)
    
    success = convert_folder_to_pptx(html_folder, output_path)
    
    if success:
        print("-" * 50)
        print("✅ 변환 완료!")
        print(f"출력 파일: {output_path}")
        print(f"파일 크기: {Path(output_path).stat().st_size:,} bytes")
        print("🎨 CSS 스타일이 적용된 디자인으로 변환되었습니다!")
    else:
        print("❌ 변환 실패!")

if __name__ == "__main__":
    main()
