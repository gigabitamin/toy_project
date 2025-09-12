#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Debug HTML to PPTX Converter
HTML 구조를 정확히 분석하고 디버깅하는 변환기
"""

import os
import sys
import requests
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup
import re
import tempfile
import shutil
from html2image import Html2Image
from PIL import Image

class DebugHTMLConverter:
    def __init__(self, html_file, output_path):
        self.html_file = html_file
        self.output_path = output_path
        self.temp_dir = None
        self.hti = Html2Image()
        
    def setup_temp_directory(self):
        """임시 디렉토리 생성"""
        self.temp_dir = Path(tempfile.mkdtemp())
        print(f"임시 디렉토리 생성: {self.temp_dir}")
        
    def cleanup_temp_directory(self):
        """임시 디렉토리 정리"""
        if self.temp_dir and self.temp_dir.exists():
            shutil.rmtree(self.temp_dir)
    
    def download_fontawesome_svg(self, icon_class, color='#2563eb'):
        """FontAwesome 아이콘 SVG 다운로드"""
        try:
            clean_class = icon_class.replace('fas ', '').replace('fab ', '').replace('far ', '').replace('fa-', '')
            print(f"아이콘 다운로드 시도: {clean_class}")
            
            paths = [
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/6.x/svgs/solid/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/6.x/svgs/brands/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/6.x/svgs/regular/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/5.x/svgs/solid/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/5.x/svgs/brands/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/5.x/svgs/regular/{clean_class}.svg",
            ]
            
            for path in paths:
                try:
                    response = requests.get(path, timeout=10)
                    if response.status_code == 200:
                        svg_content = response.text
                        svg_content = svg_content.replace('fill="currentColor"', f'fill="{color}"')
                        svg_content = svg_content.replace('fill="#000"', f'fill="{color}"')
                        svg_content = svg_content.replace('fill="black"', f'fill="{color}"')
                        
                        svg_file = self.temp_dir / f"{clean_class}.svg"
                        with open(svg_file, 'w', encoding='utf-8') as f:
                            f.write(svg_content)
                        print(f"아이콘 다운로드 성공: {clean_class}")
                        return svg_file
                except Exception as e:
                    print(f"아이콘 다운로드 실패 ({path}): {e}")
                    continue
            print(f"아이콘 다운로드 완전 실패: {clean_class}")
            return None
        except Exception as e:
            print(f"아이콘 다운로드 오류: {e}")
            return None
    
    def svg_to_png_with_html2image(self, svg_file, size=64):
        """SVG를 HTML2Image로 PNG 변환"""
        try:
            with open(svg_file, 'r', encoding='utf-8') as f:
                svg_content = f.read()
            
            html_content = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <style>
                    body {{
                        margin: 0;
                        padding: 0;
                        width: {size}px;
                        height: {size}px;
                        display: flex;
                        align-items: center;
                        justify-content: center;
                        background: transparent;
                    }}
                    svg {{
                        width: {size}px;
                        height: {size}px;
                    }}
                </style>
            </head>
            <body>
                {svg_content}
            </body>
            </html>
            """
            
            temp_html = self.temp_dir / f"{svg_file.stem}_temp.html"
            with open(temp_html, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            png_file = self.temp_dir / f"{svg_file.stem}.png"
            self.hti.screenshot(
                html_file=str(temp_html),
                save_as=f"{svg_file.stem}.png",
                size=(size, size)
            )
            
            generated_png = Path(f"{svg_file.stem}.png")
            if generated_png.exists():
                shutil.move(str(generated_png), str(png_file))
                print(f"SVG to PNG 변환 성공: {svg_file.stem}")
                return png_file
            else:
                print(f"PNG 파일 생성 실패: {svg_file.stem}")
                return None
        except Exception as e:
            print(f"SVG to PNG 변환 실패: {e}")
            return None
    
    def create_centered_text(self, slide, text, x, y, width, height, font_size=16, color='#000000', bold=False, font_family='Arial'):
        """중앙 정렬 텍스트 생성"""
        try:
            textbox = slide.shapes.add_textbox(
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            
            text_frame = textbox.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = 0
            text_frame.margin_right = 0
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            
            p = text_frame.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.name = font_family
            
            if color.startswith('#'):
                hex_color = color[1:]
                if len(hex_color) == 3:
                    hex_color = ''.join([c*2 for c in hex_color])
                p.font.color.rgb = RGBColor(
                    int(hex_color[0:2], 16),
                    int(hex_color[2:4], 16),
                    int(hex_color[4:6], 16)
                )
            
            return textbox
        except Exception as e:
            print(f"중앙 정렬 텍스트 생성 오류: {e}")
            return None
    
    def create_left_aligned_text(self, slide, text, x, y, width, height, font_size=16, color='#000000', bold=False, font_family='Arial'):
        """왼쪽 정렬 텍스트 생성"""
        try:
            textbox = slide.shapes.add_textbox(
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            
            text_frame = textbox.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = 0
            text_frame.margin_right = 0
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            
            p = text_frame.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.LEFT
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.name = font_family
            
            if color.startswith('#'):
                hex_color = color[1:]
                if len(hex_color) == 3:
                    hex_color = ''.join([c*2 for c in hex_color])
                p.font.color.rgb = RGBColor(
                    int(hex_color[0:2], 16),
                    int(hex_color[2:4], 16),
                    int(hex_color[4:6], 16)
                )
            
            return textbox
        except Exception as e:
            print(f"왼쪽 정렬 텍스트 생성 오류: {e}")
            return None
    
    def create_tech_badge(self, slide, title, x, y, width, height, icon_class=None, bg_color='#eff6ff', text_color='#1e40af', icon_color='#3b82f6'):
        """기술 스택 배지 생성"""
        try:
            print(f"기술 배지 생성: {title}")
            
            # 배지 배경 (둥근 모서리)
            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            
            # 배경 색상
            fill = badge.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(
                int(bg_color[1:3], 16),
                int(bg_color[3:5], 16),
                int(bg_color[5:7], 16)
            )
            
            # 테두리 제거
            line = badge.line
            line.fill.background()
            
            # 아이콘 추가 (왼쪽)
            if icon_class:
                print(f"아이콘 처리: {icon_class}")
                svg_file = self.download_fontawesome_svg(icon_class, icon_color)
                if svg_file:
                    png_file = self.svg_to_png_with_html2image(svg_file, 20)
                    if png_file and png_file.exists():
                        icon_x = x + 0.1
                        icon_y = y + (height - 0.25) / 2
                        icon_size = 0.25
                        
                        slide.shapes.add_picture(
                            str(png_file), 
                            Inches(icon_x), Inches(icon_y), 
                            Inches(icon_size), Inches(icon_size)
                        )
                        print(f"아이콘 추가 성공: {icon_class}")
            
            # 텍스트 추가 (아이콘 오른쪽)
            text_x = x + 0.4 if icon_class else x + 0.1
            text_width = width - 0.5 if icon_class else width - 0.2
            
            textbox = slide.shapes.add_textbox(
                Inches(text_x), Inches(y + 0.1), Inches(text_width), Inches(height - 0.2)
            )
            
            text_frame = textbox.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = 0
            text_frame.margin_right = 0
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            
            p = text_frame.paragraphs[0]
            p.text = title
            p.alignment = PP_ALIGN.LEFT
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.name = 'Arial'
            
            # 텍스트 색상
            p.font.color.rgb = RGBColor(
                int(text_color[1:3], 16),
                int(text_color[3:5], 16),
                int(text_color[5:7], 16)
            )
            
            return badge
            
        except Exception as e:
            print(f"기술 배지 생성 오류: {e}")
            return None
    
    def create_button(self, slide, text, x, y, width, height, bg_color='#3b82f6', text_color='#ffffff', icon_class=None):
        """버튼 생성"""
        try:
            print(f"버튼 생성: {text}")
            
            # 버튼 배경
            button = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            
            # 배경 색상
            fill = button.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(
                int(bg_color[1:3], 16),
                int(bg_color[3:5], 16),
                int(bg_color[5:7], 16)
            )
            
            # 테두리 제거
            line = button.line
            line.fill.background()
            
            # 아이콘 추가 (왼쪽)
            if icon_class:
                print(f"버튼 아이콘 처리: {icon_class}")
                svg_file = self.download_fontawesome_svg(icon_class, text_color)
                if svg_file:
                    png_file = self.svg_to_png_with_html2image(svg_file, 20)
                    if png_file and png_file.exists():
                        icon_x = x + 0.1
                        icon_y = y + (height - 0.25) / 2
                        icon_size = 0.25
                        
                        slide.shapes.add_picture(
                            str(png_file), 
                            Inches(icon_x), Inches(icon_y), 
                            Inches(icon_size), Inches(icon_size)
                        )
                        print(f"버튼 아이콘 추가 성공: {icon_class}")
            
            # 텍스트 추가
            text_x = x + 0.4 if icon_class else x + 0.1
            text_width = width - 0.5 if icon_class else width - 0.2
            
            textbox = slide.shapes.add_textbox(
                Inches(text_x), Inches(y + 0.1), Inches(text_width), Inches(height - 0.2)
            )
            
            text_frame = textbox.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = 0
            text_frame.margin_right = 0
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            
            p = text_frame.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.name = 'Arial'
            
            # 텍스트 색상
            p.font.color.rgb = RGBColor(
                int(text_color[1:3], 16),
                int(text_color[3:5], 16),
                int(text_color[5:7], 16)
            )
            
            return button
            
        except Exception as e:
            print(f"버튼 생성 오류: {e}")
            return None
    
    def extract_icon_class(self, element):
        """요소에서 아이콘 클래스 추출"""
        if not element:
            return None
            
        icon_elem = element.find('i')
        if icon_elem:
            classes = icon_elem.get('class', [])
            for cls in classes:
                if cls.startswith('fa-'):
                    return cls
        
        classes = element.get('class', [])
        for cls in classes:
            if cls.startswith('fa-'):
                return cls
        
        return None
    
    def parse_01_html_debug(self, soup, slide):
        """01.html 디버깅 파싱"""
        try:
            print("=== 01.html 파싱 시작 ===")
            
            # 메인 제목
            title_elem = soup.find('h1', class_='title')
            print(f"제목 요소 찾기: {title_elem}")
            if title_elem:
                print(f"제목 텍스트: {title_elem.get_text().strip()}")
                self.create_centered_text(
                    slide, title_elem.get_text().strip(),
                    2.0, 1.0, 9.0, 1.0,
                    font_size=48, color='#2563eb', bold=True
                )
            
            # 부제목
            subtitle_elem = soup.find('h2', class_='subtitle')
            print(f"부제목 요소 찾기: {subtitle_elem}")
            if subtitle_elem:
                print(f"부제목 텍스트: {subtitle_elem.get_text().strip()}")
                self.create_centered_text(
                    slide, subtitle_elem.get_text().strip(),
                    2.0, 2.2, 9.0, 0.6,
                    font_size=32, color='#1e40af', bold=True
                )
            
            # Project Info Section 찾기
            project_info = soup.find('div', class_='text-center mb-16')
            print(f"Project Info Section 찾기: {project_info}")
            
            if project_info:
                # 개발 기간 섹션
                period_section = project_info.find('div', class_='mb-10')
                print(f"개발 기간 섹션 찾기: {period_section}")
                
                if period_section:
                    # 개발 기간 제목
                    period_title = period_section.find('p', class_='text-xl mb-2 text-gray-600')
                    print(f"개발 기간 제목: {period_title}")
                    if period_title:
                        print(f"개발 기간 제목 텍스트: {period_title.get_text().strip()}")
                        self.create_centered_text(
                            slide, period_title.get_text().strip(),
                            2.0, 3.0, 9.0, 0.4,
                            font_size=20, color='#6b7280', bold=False
                        )
                    
                    # 개발 기간 날짜
                    period_date = period_section.find('p', class_='text-2xl font-medium')
                    print(f"개발 기간 날짜: {period_date}")
                    if period_date:
                        print(f"개발 기간 날짜 텍스트: {period_date.get_text().strip()}")
                        self.create_centered_text(
                            slide, period_date.get_text().strip(),
                            2.0, 3.5, 9.0, 0.5,
                            font_size=24, color='#1f2937', bold=True
                        )
                
                # 기술 스택 섹션들 찾기
                tech_sections = project_info.find_all('div', class_='mb-10')
                print(f"기술 스택 섹션들 개수: {len(tech_sections)}")
                
                for i, tech_section in enumerate(tech_sections):
                    print(f"기술 스택 섹션 {i}: {tech_section}")
                    
                    # 기술 스택 제목 찾기
                    tech_title = tech_section.find('p', class_='text-xl mb-4 text-gray-600')
                    print(f"기술 스택 제목: {tech_title}")
                    
                    if tech_title and '기술 스택' in tech_title.get_text():
                        print(f"기술 스택 제목 텍스트: {tech_title.get_text().strip()}")
                        self.create_centered_text(
                            slide, tech_title.get_text().strip(),
                            2.0, 4.2, 9.0, 0.4,
                            font_size=20, color='#6b7280', bold=False
                        )
                        
                        # 기술 스택 배지들 찾기
                        tech_badges_container = tech_section.find('div', class_='flex flex-wrap justify-center')
                        print(f"기술 스택 배지 컨테이너: {tech_badges_container}")
                        
                        if tech_badges_container:
                            tech_badges = tech_badges_container.find_all('div', class_='tech-stack')
                            print(f"기술 스택 배지들 개수: {len(tech_badges)}")
                            
                            y_pos = 4.8
                            
                            for j, badge in enumerate(tech_badges[:5]):  # 최대 5개
                                text = badge.get_text().strip()
                                icon_class = self.extract_icon_class(badge)
                                print(f"기술 배지 {j}: {text}, 아이콘: {icon_class}")
                                
                                x_pos = 1.0 + (j % 3) * 3.5
                                y_badge = y_pos + (j // 3) * 0.8
                                
                                self.create_tech_badge(
                                    slide, text,
                                    x_pos, y_badge, 3.0, 0.6,
                                    icon_class, '#eff6ff', '#1e40af', '#3b82f6'
                                )
            
            # 링크 버튼들
            link_section = soup.find('div', class_='flex justify-center space-x-8 mt-4')
            print(f"링크 섹션: {link_section}")
            
            if link_section:
                buttons = link_section.find_all('a', class_='link-button')
                print(f"버튼들 개수: {len(buttons)}")
                
                y_pos = 6.5
                
                for i, button in enumerate(buttons[:2]):
                    text = button.get_text().strip()
                    icon_class = self.extract_icon_class(button)
                    print(f"버튼 {i}: {text}, 아이콘: {icon_class}")
                    
                    x_pos = 2.0 + i * 4.5
                    
                    if 'github' in text.lower():
                        self.create_button(
                            slide, text, x_pos, y_pos, 3.5, 0.6,
                            '#1f2937', '#ffffff', icon_class
                        )
                    else:
                        self.create_button(
                            slide, text, x_pos, y_pos, 3.5, 0.6,
                            '#2563eb', '#ffffff', icon_class
                        )
            
            # 하단 날짜
            footer = soup.find('div', class_='absolute bottom-8 right-8 text-gray-500')
            print(f"하단 날짜: {footer}")
            if footer:
                date_text = footer.find('p')
                if date_text:
                    print(f"날짜 텍스트: {date_text.get_text().strip()}")
                    self.create_left_aligned_text(
                        slide, date_text.get_text().strip(),
                        10.0, 6.8, 2.0, 0.3,
                        font_size=12, color='#6b7280', bold=False
                    )
            
            print("=== 01.html 파싱 완료 ===")
                    
        except Exception as e:
            print(f"01.html 파싱 오류: {e}")
            import traceback
            traceback.print_exc()
    
    def parse_02_html_debug(self, soup, slide):
        """02.html 디버깅 파싱"""
        try:
            print("=== 02.html 파싱 시작 ===")
            
            # 섹션 제목
            title_elem = soup.find('h1', class_='section-title')
            print(f"섹션 제목 요소 찾기: {title_elem}")
            if title_elem:
                print(f"섹션 제목 텍스트: {title_elem.get_text().strip()}")
                self.create_left_aligned_text(
                    slide, title_elem.get_text().strip(),
                    1.0, 0.5, 11.0, 0.8,
                    font_size=36, color='#2563eb', bold=True
                )
            
            # 구분선
            divider = soup.find('div', class_='w-24 h-1 bg-blue-500 mb-8')
            print(f"구분선 요소 찾기: {divider}")
            if divider:
                print("구분선 추가")
                # 구분선 그리기
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(1.0), Inches(1.8), Inches(1.5), Inches(0.05)
                )
                fill = line.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(59, 130, 246)  # #3b82f6
                line.line.fill.background()
            
            # 배경 섹션
            background_section = soup.find('div', class_='flex items-start')
            print(f"배경 섹션 찾기: {background_section}")
            
            if background_section:
                # 아이콘
                icon_elem = background_section.find('i', class_='fas fa-history')
                print(f"배경 아이콘: {icon_elem}")
                if icon_elem:
                    print("배경 아이콘 원형 생성")
                    # 아이콘 원형 배경
                    circle = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        Inches(1.0), Inches(2.2), Inches(0.8), Inches(0.8)
                    )
                    fill = circle.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(219, 234, 254)  # #dbeafe
                    
                    # 아이콘 이미지 추가
                    svg_file = self.download_fontawesome_svg('fa-history', '#2563eb')
                    if svg_file:
                        png_file = self.svg_to_png_with_html2image(svg_file, 32)
                        if png_file and png_file.exists():
                            slide.shapes.add_picture(
                                str(png_file), 
                                Inches(1.2), Inches(2.4), 
                                Inches(0.4), Inches(0.4)
                            )
                            print("배경 아이콘 추가 성공")
                
                # 제목과 내용
                title_elem = background_section.find('h2', class_='text-2xl font-bold mb-3 text-gray-800')
                content_elem = background_section.find('p', class_='text-lg text-gray-600 leading-relaxed')
                
                print(f"배경 제목: {title_elem}")
                print(f"배경 내용: {content_elem}")
                
                if title_elem:
                    print(f"배경 제목 텍스트: {title_elem.get_text().strip()}")
                    self.create_left_aligned_text(
                        slide, title_elem.get_text().strip(),
                        2.0, 2.2, 9.0, 0.5,
                        font_size=24, color='#1f2937', bold=True
                    )
                
                if content_elem:
                    print(f"배경 내용 텍스트: {content_elem.get_text().strip()}")
                    self.create_left_aligned_text(
                        slide, content_elem.get_text().strip(),
                        2.0, 2.8, 9.0, 1.0,
                        font_size=18, color='#6b7280', bold=False
                    )
            
            # 목적 섹션
            purpose_sections = soup.find_all('div', class_='flex items-start')
            print(f"목적 섹션들 개수: {len(purpose_sections)}")
            
            if len(purpose_sections) > 1:
                purpose_section = purpose_sections[1]
                print(f"목적 섹션: {purpose_section}")
                
                # 아이콘
                icon_elem = purpose_section.find('i', class_='fas fa-bullseye')
                print(f"목적 아이콘: {icon_elem}")
                if icon_elem:
                    print("목적 아이콘 원형 생성")
                    circle = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        Inches(1.0), Inches(4.0), Inches(0.8), Inches(0.8)
                    )
                    fill = circle.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(219, 234, 254)  # #dbeafe
                    
                    svg_file = self.download_fontawesome_svg('fa-bullseye', '#2563eb')
                    if svg_file:
                        png_file = self.svg_to_png_with_html2image(svg_file, 32)
                        if png_file and png_file.exists():
                            slide.shapes.add_picture(
                                str(png_file), 
                                Inches(1.2), Inches(4.2), 
                                Inches(0.4), Inches(0.4)
                            )
                            print("목적 아이콘 추가 성공")
                
                # 제목과 내용
                title_elem = purpose_section.find('h2', class_='text-2xl font-bold mb-3 text-gray-800')
                content_elem = purpose_section.find('p', class_='text-lg text-gray-600 leading-relaxed')
                
                print(f"목적 제목: {title_elem}")
                print(f"목적 내용: {content_elem}")
                
                if title_elem:
                    print(f"목적 제목 텍스트: {title_elem.get_text().strip()}")
                    self.create_left_aligned_text(
                        slide, title_elem.get_text().strip(),
                        2.0, 4.0, 9.0, 0.5,
                        font_size=24, color='#1f2937', bold=True
                    )
                
                if content_elem:
                    print(f"목적 내용 텍스트: {content_elem.get_text().strip()}")
                    self.create_left_aligned_text(
                        slide, content_elem.get_text().strip(),
                        2.0, 4.6, 9.0, 1.0,
                        font_size=18, color='#6b7280', bold=False
                    )
            
            # 주요 특징 섹션
            if len(purpose_sections) > 2:
                features_section = purpose_sections[2]
                print(f"주요 특징 섹션: {features_section}")
                
                # 아이콘
                icon_elem = features_section.find('i', class_='fas fa-star')
                print(f"주요 특징 아이콘: {icon_elem}")
                if icon_elem:
                    print("주요 특징 아이콘 원형 생성")
                    circle = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        Inches(1.0), Inches(5.8), Inches(0.8), Inches(0.8)
                    )
                    fill = circle.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(219, 234, 254)  # #dbeafe
                    
                    svg_file = self.download_fontawesome_svg('fa-star', '#2563eb')
                    if svg_file:
                        png_file = self.svg_to_png_with_html2image(svg_file, 32)
                        if png_file and png_file.exists():
                            slide.shapes.add_picture(
                                str(png_file), 
                                Inches(1.2), Inches(6.0), 
                                Inches(0.4), Inches(0.4)
                            )
                            print("주요 특징 아이콘 추가 성공")
                
                # 제목
                title_elem = features_section.find('h2', class_='text-2xl font-bold mb-4 text-gray-800')
                print(f"주요 특징 제목: {title_elem}")
                if title_elem:
                    print(f"주요 특징 제목 텍스트: {title_elem.get_text().strip()}")
                    self.create_left_aligned_text(
                        slide, title_elem.get_text().strip(),
                        2.0, 5.8, 9.0, 0.5,
                        font_size=24, color='#1f2937', bold=True
                    )
                
                # 기능 카드들
                feature_cards = features_section.find_all('div', class_='feature-card')
                print(f"기능 카드들 개수: {len(feature_cards)}")
                
                y_pos = 6.5
                
                for i, card in enumerate(feature_cards[:4]):  # 최대 4개
                    title_elem = card.find('h3', class_='font-bold text-lg mb-1')
                    content_elem = card.find('p', class_='text-gray-600')
                    icon_elem = card.find('i')
                    
                    title = title_elem.get_text().strip() if title_elem else ""
                    content = content_elem.get_text().strip() if content_elem else ""
                    icon_class = self.extract_icon_class(card)
                    
                    print(f"기능 카드 {i}: {title}, {content}, 아이콘: {icon_class}")
                    
                    x_pos = 0.5 + (i % 2) * 6.0
                    y_card = y_pos + (i // 2) * 1.5
                    
                    # 카드 배경
                    card_shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(x_pos), Inches(y_card), Inches(5.5), Inches(1.2)
                    )
                    
                    fill = card_shape.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(249, 250, 251)  # bg-gray-50
                    
                    line = card_shape.line
                    line.color.rgb = RGBColor(229, 231, 235)  # border-gray-200
                    line.width = Pt(1)
                    
                    # 아이콘 추가
                    if icon_class:
                        svg_file = self.download_fontawesome_svg(icon_class, '#3b82f6')
                        if svg_file:
                            png_file = self.svg_to_png_with_html2image(svg_file, 24)
                            if png_file and png_file.exists():
                                slide.shapes.add_picture(
                                    str(png_file), 
                                    Inches(x_pos + 0.2), Inches(y_card + 0.2), 
                                    Inches(0.3), Inches(0.3)
                                )
                    
                    # 제목 텍스트
                    title_x = x_pos + 0.6 if icon_class else x_pos + 0.2
                    title_width = 5.5 - 0.8 if icon_class else 5.5 - 0.4
                    
                    title_box = slide.shapes.add_textbox(
                        Inches(title_x), Inches(y_card + 0.2), Inches(title_width), Inches(0.4)
                    )
                    title_frame = title_box.text_frame
                    title_frame.clear()
                    p1 = title_frame.paragraphs[0]
                    p1.text = title
                    p1.alignment = PP_ALIGN.LEFT
                    p1.font.size = Pt(16)
                    p1.font.bold = True
                    p1.font.name = 'Arial'
                    p1.font.color.rgb = RGBColor(31, 41, 55)
                    
                    # 설명 텍스트
                    desc_box = slide.shapes.add_textbox(
                        Inches(title_x), Inches(y_card + 0.7), Inches(title_width), Inches(0.5)
                    )
                    desc_frame = desc_box.text_frame
                    desc_frame.clear()
                    p2 = desc_frame.paragraphs[0]
                    p2.text = content
                    p2.alignment = PP_ALIGN.LEFT
                    p2.font.size = Pt(12)
                    p2.font.name = 'Arial'
                    p2.font.color.rgb = RGBColor(107, 114, 128)
            
            print("=== 02.html 파싱 완료 ===")
                    
        except Exception as e:
            print(f"02.html 파싱 오류: {e}")
            import traceback
            traceback.print_exc()

    def parse_03_html_debug(self, soup, slide):
        """03.html 디버깅 파싱"""
        try:
            print("=== 03.html 파싱 시작 ===")
            
            # 헤더 섹션
            header_section = soup.find('div', class_='mb-8')
            print(f"헤더 섹션 찾기: {header_section}")
            
            if header_section:
                # 제목
                title_elem = header_section.find('h1', class_='section-title')
                print(f"제목 요소: {title_elem}")
                if title_elem:
                    print(f"제목 텍스트: {title_elem.get_text().strip()}")
                    self.create_left_aligned_text(
                        slide, title_elem.get_text().strip(),
                        1.0, 0.5, 11.0, 0.8,
                        font_size=32, color='#2563eb', bold=True
                    )
                
                # 구분선
                divider = header_section.find('div', class_='w-20 h-1 bg-blue-500')
                print(f"구분선 요소: {divider}")
                if divider:
                    print("구분선 추가")
                    line = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(1.0), Inches(1.5), Inches(1.25), Inches(0.05)
                    )
                    fill = line.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(59, 130, 246)  # #3b82f6
                    line.line.fill.background()
            
            # 메인 콘텐츠 섹션
            content_section = soup.find('div', class_='grid grid-cols-2 gap-8')
            print(f"메인 콘텐츠 섹션: {content_section}")
            
            if content_section:
                # Frontend & Backend 섹션
                frontend_section = content_section.find('div')
                print(f"Frontend 섹션: {frontend_section}")
                
                if frontend_section:
                    # 섹션 제목
                    section_title = frontend_section.find('h2', class_='text-xl font-bold mb-4')
                    print(f"Frontend 섹션 제목: {section_title}")
                    if section_title:
                        print(f"Frontend 섹션 제목 텍스트: {section_title.get_text().strip()}")
                        self.create_left_aligned_text(
                            slide, section_title.get_text().strip(),
                            1.0, 2.0, 5.0, 0.5,
                            font_size=20, color='#2563eb', bold=True
                        )
                    
                    # 기술 카드들
                    tech_cards = frontend_section.find_all('div', class_='tech-card')
                    print(f"기술 카드들 개수: {len(tech_cards)}")
                    
                    y_pos = 2.8
                    for i, card in enumerate(tech_cards[:6]):  # 최대 6개
                        title_elem = card.find('h3', class_='font-bold text-sm')
                        desc_elem = card.find('p', class_='text-gray-600 text-xs')
                        icon_elem = card.find('i')
                        
                        title = title_elem.get_text().strip() if title_elem else ""
                        desc = desc_elem.get_text().strip() if desc_elem else ""
                        icon_class = self.extract_icon_class(card)
                        
                        print(f"기술 카드 {i}: {title}, {desc}, 아이콘: {icon_class}")
                        
                        x_pos = 1.0 + (i % 2) * 2.5
                        y_card = y_pos + (i // 2) * 1.2
                        
                        # 카드 배경
                        card_shape = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x_pos), Inches(y_card), Inches(2.2), Inches(1.0)
                        )
                        
                        fill = card_shape.fill
                        fill.solid()
                        if 'frontend-card' in card.get('class', []):
                            fill.fore_color.rgb = RGBColor(239, 246, 255)  # bg-blue-50
                        elif 'backend-card' in card.get('class', []):
                            fill.fore_color.rgb = RGBColor(240, 253, 244)  # bg-green-50
                        
                        line = card_shape.line
                        line.color.rgb = RGBColor(229, 231, 235)  # border-gray-200
                        line.width = Pt(1)
                        
                        # 아이콘 추가
                        if icon_class:
                            svg_file = self.download_fontawesome_svg(icon_class, '#3b82f6')
                            if svg_file:
                                png_file = self.svg_to_png_with_html2image(svg_file, 20)
                                if png_file and png_file.exists():
                                    slide.shapes.add_picture(
                                        str(png_file), 
                                        Inches(x_pos + 0.1), Inches(y_card + 0.1), 
                                        Inches(0.3), Inches(0.3)
                                    )
                        
                        # 제목 텍스트
                        title_x = x_pos + 0.5 if icon_class else x_pos + 0.1
                        title_width = 2.2 - 0.6 if icon_class else 2.2 - 0.2
                        
                        title_box = slide.shapes.add_textbox(
                            Inches(title_x), Inches(y_card + 0.1), Inches(title_width), Inches(0.3)
                        )
                        title_frame = title_box.text_frame
                        title_frame.clear()
                        p1 = title_frame.paragraphs[0]
                        p1.text = title
                        p1.alignment = PP_ALIGN.LEFT
                        p1.font.size = Pt(12)
                        p1.font.bold = True
                        p1.font.name = 'Arial'
                        p1.font.color.rgb = RGBColor(31, 41, 55)
                        
                        # 설명 텍스트
                        desc_box = slide.shapes.add_textbox(
                            Inches(title_x), Inches(y_card + 0.5), Inches(title_width), Inches(0.4)
                        )
                        desc_frame = desc_box.text_frame
                        desc_frame.clear()
                        p2 = desc_frame.paragraphs[0]
                        p2.text = desc
                        p2.alignment = PP_ALIGN.LEFT
                        p2.font.size = Pt(10)
                        p2.font.name = 'Arial'
                        p2.font.color.rgb = RGBColor(107, 114, 128)
                
                # 학습 성과 섹션
                learning_section = frontend_section.find_next_sibling('div')
                print(f"학습 성과 섹션: {learning_section}")
                
                if learning_section:
                    # 섹션 제목
                    section_title = learning_section.find('h2', class_='text-xl font-bold mb-4')
                    print(f"학습 성과 섹션 제목: {section_title}")
                    if section_title:
                        print(f"학습 성과 섹션 제목 텍스트: {section_title.get_text().strip()}")
                        self.create_left_aligned_text(
                            slide, section_title.get_text().strip(),
                            6.5, 2.0, 5.0, 0.5,
                            font_size=20, color='#8b5cf6', bold=True
                        )
                    
                    # 학습 성과 카드들
                    learning_cards = learning_section.find_all('div', class_='mb-4')
                    print(f"학습 성과 카드들 개수: {len(learning_cards)}")
                    
                    y_pos = 2.8
                    for i, card in enumerate(learning_cards[:4]):  # 최대 4개
                        title_elem = card.find('h3', class_='text-lg font-bold mb-2')
                        print(f"학습 성과 카드 {i} 제목: {title_elem}")
                        
                        if title_elem:
                            title = title_elem.get_text().strip()
                            print(f"학습 성과 카드 {i} 제목 텍스트: {title}")
                            
                            # 제목 추가
                            self.create_left_aligned_text(
                                slide, title,
                                6.5, y_pos, 5.0, 0.4,
                                font_size=16, color='#2563eb', bold=True
                            )
                            
                            # 하위 카드들
                            sub_cards = card.find_all('div', class_='tech-card')
                            print(f"하위 카드들 개수: {len(sub_cards)}")
                            
                            for j, sub_card in enumerate(sub_cards[:3]):  # 최대 3개
                                sub_title_elem = sub_card.find('h4', class_='font-bold text-xs')
                                sub_desc_elem = sub_card.find('p', class_='text-gray-600 text-xs')
                                sub_icon_elem = sub_card.find('i')
                                
                                sub_title = sub_title_elem.get_text().strip() if sub_title_elem else ""
                                sub_desc = sub_desc_elem.get_text().strip() if sub_desc_elem else ""
                                sub_icon_class = self.extract_icon_class(sub_card)
                                
                                print(f"하위 카드 {j}: {sub_title}, {sub_desc}, 아이콘: {sub_icon_class}")
                                
                                sub_x = 6.5 + j * 1.8
                                sub_y = y_pos + 0.5
                                
                                # 하위 카드 배경
                                sub_card_shape = slide.shapes.add_shape(
                                    MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(sub_x), Inches(sub_y), Inches(1.6), Inches(0.8)
                                )
                                
                                fill = sub_card_shape.fill
                                fill.solid()
                                fill.fore_color.rgb = RGBColor(249, 250, 251)  # bg-gray-50
                                
                                line = sub_card_shape.line
                                line.color.rgb = RGBColor(229, 231, 235)
                                line.width = Pt(1)
                                
                                # 하위 카드 아이콘
                                if sub_icon_class:
                                    svg_file = self.download_fontawesome_svg(sub_icon_class, '#3b82f6')
                                    if svg_file:
                                        png_file = self.svg_to_png_with_html2image(svg_file, 16)
                                        if png_file and png_file.exists():
                                            slide.shapes.add_picture(
                                                str(png_file), 
                                                Inches(sub_x + 0.05), Inches(sub_y + 0.05), 
                                                Inches(0.2), Inches(0.2)
                                            )
                                
                                # 하위 카드 제목
                                sub_title_x = sub_x + 0.3 if sub_icon_class else sub_x + 0.05
                                sub_title_width = 1.6 - 0.35 if sub_icon_class else 1.6 - 0.1
                                
                                sub_title_box = slide.shapes.add_textbox(
                                    Inches(sub_title_x), Inches(sub_y + 0.05), Inches(sub_title_width), Inches(0.25)
                                )
                                sub_title_frame = sub_title_box.text_frame
                                sub_title_frame.clear()
                                p1 = sub_title_frame.paragraphs[0]
                                p1.text = sub_title
                                p1.alignment = PP_ALIGN.LEFT
                                p1.font.size = Pt(9)
                                p1.font.bold = True
                                p1.font.name = 'Arial'
                                p1.font.color.rgb = RGBColor(31, 41, 55)
                                
                                # 하위 카드 설명
                                sub_desc_box = slide.shapes.add_textbox(
                                    Inches(sub_title_x), Inches(sub_y + 0.35), Inches(sub_title_width), Inches(0.4)
                                )
                                sub_desc_frame = sub_desc_box.text_frame
                                sub_desc_frame.clear()
                                p2 = sub_desc_frame.paragraphs[0]
                                p2.text = sub_desc
                                p2.alignment = PP_ALIGN.LEFT
                                p2.font.size = Pt(8)
                                p2.font.name = 'Arial'
                                p2.font.color.rgb = RGBColor(107, 114, 128)
                            
                            y_pos += 1.8
            
            # 배포 & 호스팅 섹션
            deploy_section = soup.find('div', class_='mt-8')
            print(f"배포 섹션: {deploy_section}")
            
            if deploy_section:
                # 섹션 제목
                section_title = deploy_section.find('h2', class_='text-xl font-bold mb-6')
                print(f"배포 섹션 제목: {section_title}")
                if section_title:
                    print(f"배포 섹션 제목 텍스트: {section_title.get_text().strip()}")
                    self.create_left_aligned_text(
                        slide, section_title.get_text().strip(),
                        1.0, 5.5, 11.0, 0.5,
                        font_size=20, color='#8b5cf6', bold=True
                    )
                
                # 배포 플로우
                flow_container = deploy_section.find('div', class_='bg-gray-50 rounded-lg p-6')
                print(f"배포 플로우 컨테이너: {flow_container}")
                
                if flow_container:
                    flow_items = flow_container.find_all('div', class_='flex flex-col items-center')
                    print(f"배포 플로우 아이템들 개수: {len(flow_items)}")
                    
                    x_pos = 2.0
                    for i, item in enumerate(flow_items):
                        icon_elem = item.find('i')
                        text_elem = item.find('p', class_='font-medium text-sm')
                        
                        if icon_elem and text_elem:
                            icon_class = self.extract_icon_class(item)
                            text = text_elem.get_text().strip()
                            
                            print(f"배포 플로우 {i}: {text}, 아이콘: {icon_class}")
                            
                            # 아이콘 원형 배경
                            circle = slide.shapes.add_shape(
                                MSO_SHAPE.OVAL,
                                Inches(x_pos), Inches(6.2), Inches(0.8), Inches(0.8)
                            )
                            fill = circle.fill
                            fill.solid()
                            fill.fore_color.rgb = RGBColor(239, 246, 255)  # bg-blue-100
                            
                            # 아이콘 이미지
                            if icon_class:
                                svg_file = self.download_fontawesome_svg(icon_class, '#2563eb')
                                if svg_file:
                                    png_file = self.svg_to_png_with_html2image(svg_file, 24)
                                    if png_file and png_file.exists():
                                        slide.shapes.add_picture(
                                            str(png_file), 
                                            Inches(x_pos + 0.2), Inches(6.4), 
                                            Inches(0.4), Inches(0.4)
                                        )
                            
                            # 텍스트
                            text_box = slide.shapes.add_textbox(
                                Inches(x_pos - 0.2), Inches(7.2), Inches(1.2), Inches(0.3)
                            )
                            text_frame = text_box.text_frame
                            text_frame.clear()
                            p = text_frame.paragraphs[0]
                            p.text = text
                            p.alignment = PP_ALIGN.CENTER
                            p.font.size = Pt(10)
                            p.font.bold = True
                            p.font.name = 'Arial'
                            p.font.color.rgb = RGBColor(31, 41, 55)
                            
                            x_pos += 2.0
                        
                        # 화살표 추가 (마지막이 아닌 경우)
                        if i < len(flow_items) - 1:
                            arrow_box = slide.shapes.add_textbox(
                                Inches(x_pos - 0.5), Inches(6.5), Inches(0.3), Inches(0.3)
                            )
                            arrow_frame = arrow_box.text_frame
                            arrow_frame.clear()
                            p = arrow_frame.paragraphs[0]
                            p.text = "→"
                            p.alignment = PP_ALIGN.CENTER
                            p.font.size = Pt(16)
                            p.font.name = 'Arial'
                            p.font.color.rgb = RGBColor(107, 114, 128)
            
            # 푸터
            footer = soup.find('div', class_='absolute bottom-8 right-8')
            print(f"푸터: {footer}")
            if footer:
                footer_text = footer.find('p')
                if footer_text:
                    print(f"푸터 텍스트: {footer_text.get_text().strip()}")
                    self.create_left_aligned_text(
                        slide, footer_text.get_text().strip(),
                        8.0, 6.8, 4.0, 0.3,
                        font_size=10, color='#9ca3af', bold=False
                    )
            
            print("=== 03.html 파싱 완료 ===")
                    
        except Exception as e:
            print(f"03.html 파싱 오류: {e}")
            import traceback
            traceback.print_exc()

    def parse_html_debug(self, soup, slide):
        """디버깅 HTML 파싱"""
        try:
            filename = Path(self.html_file).name
            print(f"파일명: {filename}")
            
            if filename == '01.html':
                self.parse_01_html_debug(soup, slide)
            elif filename == '02.html':
                self.parse_02_html_debug(soup, slide)
            elif filename == '03.html':
                self.parse_03_html_debug(soup, slide)
            else:
                # 기타 파일들에 대한 기본 파싱
                print("기본 파싱 실행")
                
        except Exception as e:
            print(f"HTML 파싱 오류: {e}")
            import traceback
            traceback.print_exc()
    
    def convert(self):
        """HTML을 PPTX로 변환"""
        try:
            self.setup_temp_directory()
            
            with open(self.html_file, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            prs = Presentation()
            prs.slide_width = Inches(13.33)  # 16:9 비율
            prs.slide_height = Inches(7.5)
            
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            
            self.parse_html_debug(soup, slide)
            
            prs.save(self.output_path)
            print(f"✅ 변환 완료: {Path(self.html_file).name}")
            
            return True
            
        except Exception as e:
            print(f"❌ 변환 실패: {e}")
            import traceback
            traceback.print_exc()
            return False
        
        finally:
            # 디버깅을 위해 임시 디렉토리 정리하지 않음
            print(f"임시 디렉토리 유지: {self.temp_dir}")

def convert_folder_to_pptx(html_folder, output_path):
    """폴더 내 모든 HTML 파일을 하나의 PPTX로 변환"""
    try:
        html_folder = Path(html_folder)
        if not html_folder.exists():
            print(f"HTML 폴더가 존재하지 않습니다: {html_folder}")
            return False
        
        html_files = list(html_folder.glob("*.html"))
        if not html_files:
            print(f"HTML 파일이 없습니다: {html_folder}")
            return False
        
        print(f"발견된 HTML 파일 {len(html_files)}개:")
        for html_file in html_files:
            print(f"  - {html_file.name}")
        
        prs = Presentation()
        prs.slide_width = Inches(13.33)  # 16:9 비율
        prs.slide_height = Inches(7.5)
        
        for i, html_file in enumerate(html_files):
            print(f"\n--- {html_file.name} 변환 중 ({i+1}/{len(html_files)}) ---")
            
            try:
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                
                converter = DebugHTMLConverter(str(html_file), "")
                converter.setup_temp_directory()
                
                with open(html_file, 'r', encoding='utf-8') as f:
                    html_content = f.read()
                
                soup = BeautifulSoup(html_content, 'html.parser')
                converter.parse_html_debug(soup, slide)
                # 임시 디렉토리 정리하지 않음
                
                print(f"✅ {html_file.name} 변환 완료")
                
            except Exception as e:
                print(f"❌ {html_file.name} 변환 실패: {e}")
                import traceback
                traceback.print_exc()
                continue
        
        prs.save(output_path)
        print(f"\n✅ 모든 HTML 파일이 하나의 PPTX로 변환 완료!")
        print(f"출력 파일: {output_path}")
        return True
        
    except Exception as e:
        print(f"폴더 변환 오류: {e}")
        import traceback
        traceback.print_exc()
        return False

def main():
    html_folder = r"C:\Project\gigabitamin\genspark\dcs_site\html"
    output_path = r"C:\Project\gigabitamin\genspark\dcs_site\html\debug_all_pages.pptx"
    
    print("Debug HTML to Editable PPTX 변환기")
    print("=" * 50)
    print(f"HTML 폴더: {html_folder}")
    print(f"출력 파일: {output_path}")
    print("-" * 50)
    
    success = convert_folder_to_pptx(html_folder, output_path)
    
    if success:
        print("-" * 50)
        print("✅ 변환 완료!")
        print(f"출력 파일: {output_path}")
        print(f"파일 크기: {Path(output_path).stat().st_size:,} bytes")
        print("🔍 디버깅 정보가 출력되었습니다!")
    else:
        print("❌ 변환 실패!")

if __name__ == "__main__":
    main()
