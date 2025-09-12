#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
HTML to Editable PPTX Converter V3
HTML 파일을 파싱하여 텍스트, 도형, 이미지를 편집 가능한 PPTX 객체로 변환하는 스크립트
- 01.html (메인 페이지) 및 02.html (프로젝트 개요) 지원
- 레이아웃 최적화
"""

import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup
import re
import requests
from urllib.parse import urljoin, urlparse
import tempfile
import shutil

class HTMLEditablePPTXConverterV3:
    def __init__(self, html_file, output_path):
        self.html_file = Path(html_file)
        self.output_path = Path(output_path)
        self.temp_dir = None
        
    def setup_temp_directory(self):
        """임시 디렉토리 설정"""
        self.temp_dir = Path(tempfile.mkdtemp())
        print(f"임시 디렉토리 생성: {self.temp_dir}")
        
    def cleanup_temp_directory(self):
        """임시 디렉토리 정리"""
        if self.temp_dir and self.temp_dir.exists():
            shutil.rmtree(self.temp_dir)
            print("임시 디렉토리 정리 완료")
    
    def parse_css_color(self, color_str):
        """CSS 색상 문자열을 RGBColor로 변환"""
        if not color_str:
            return None
            
        # #RRGGBB 형식
        if color_str.startswith('#'):
            hex_color = color_str[1:]
            if len(hex_color) == 6:
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
                return RGBColor(r, g, b)
        
        # rgb(r, g, b) 형식
        rgb_match = re.match(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', color_str)
        if rgb_match:
            r, g, b = map(int, rgb_match.groups())
            return RGBColor(r, g, b)
        
        # 색상 이름 매핑
        color_map = {
            'blue': RGBColor(37, 99, 235),
            'red': RGBColor(239, 68, 68),
            'green': RGBColor(34, 197, 94),
            'yellow': RGBColor(234, 179, 8),
            'purple': RGBColor(147, 51, 234),
            'gray': RGBColor(107, 114, 128),
            'black': RGBColor(0, 0, 0),
            'white': RGBColor(255, 255, 255)
        }
        
        return color_map.get(color_str.lower())
    
    def parse_font_size(self, font_size_str):
        """CSS 폰트 크기를 Pt로 변환"""
        if not font_size_str:
            return 12
        
        # px 단위 제거
        if font_size_str.endswith('px'):
            return int(float(font_size_str[:-2]))
        elif font_size_str.endswith('rem'):
            return int(float(font_size_str[:-3]) * 16)  # 1rem = 16px 가정
        elif font_size_str.endswith('em'):
            return int(float(font_size_str[:-2]) * 16)  # 1em = 16px 가정
        else:
            try:
                return int(float(font_size_str))
            except:
                return 12
    
    def get_text_alignment(self, text_align):
        """CSS text-align을 PPTX 정렬로 변환"""
        align_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        return align_map.get(text_align, PP_ALIGN.LEFT)
    
    def create_text_box(self, slide, text, x, y, width, height, styles=None):
        """텍스트 박스 생성"""
        try:
            # 텍스트 박스 추가
            textbox = slide.shapes.add_textbox(
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            
            # 텍스트 프레임 설정
            text_frame = textbox.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.05)
            text_frame.margin_bottom = Inches(0.05)
            
            # 단락 생성
            p = text_frame.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.LEFT
            
            # 스타일 적용
            if styles:
                font = p.font
                
                # 폰트 크기
                if 'font-size' in styles:
                    font.size = Pt(self.parse_font_size(styles['font-size']))
                
                # 폰트 색상
                if 'color' in styles:
                    color = self.parse_css_color(styles['color'])
                    if color:
                        font.color.rgb = color
                
                # 폰트 굵기
                if 'font-weight' in styles:
                    if styles['font-weight'] in ['bold', '700', '800', '900']:
                        font.bold = True
                
                # 정렬
                if 'text-align' in styles:
                    p.alignment = self.get_text_alignment(styles['text-align'])
            
            return textbox
            
        except Exception as e:
            print(f"텍스트 박스 생성 오류: {e}")
            return None
    
    def create_tech_stack_box(self, slide, tech_text, x, y):
        """기술 스택 박스 생성"""
        try:
            # 배경 박스 생성
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(2.2), Inches(0.6)
            )
            
            # 배경 색상 설정 (blue-50)
            fill = box.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(239, 246, 255)
            
            # 테두리 설정
            line = box.line
            line.color.rgb = RGBColor(219, 234, 254)
            
            # 텍스트 추가
            text_frame = box.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            p = text_frame.paragraphs[0]
            p.text = tech_text
            p.alignment = PP_ALIGN.CENTER
            
            # 폰트 설정
            font = p.font
            font.size = Pt(12)
            font.bold = True
            font.color.rgb = RGBColor(30, 64, 175)  # blue-800
            
            return box
            
        except Exception as e:
            print(f"기술 스택 박스 생성 오류: {e}")
            return None
    
    def create_button(self, slide, text, x, y, width, height, bg_color, text_color):
        """버튼 생성"""
        try:
            # 버튼 박스 생성
            button = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            
            # 배경 색상 설정
            fill = button.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
            
            # 텍스트 추가
            text_frame = button.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            p = text_frame.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.CENTER
            
            # 폰트 설정
            font = p.font
            font.size = Pt(14)
            font.bold = True
            font.color.rgb = text_color
            
            return button
            
        except Exception as e:
            print(f"버튼 생성 오류: {e}")
            return None
    
    def create_feature_card(self, slide, title, description, x, y, width, height):
        """기능 카드 생성"""
        try:
            # 카드 배경 생성
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            
            # 배경 색상 설정 (gray-50)
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(249, 250, 251)
            
            # 테두리 설정
            line = card.line
            line.color.rgb = RGBColor(229, 231, 235)
            
            # 텍스트 추가
            text_frame = card.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)
            text_frame.margin_top = Inches(0.2)
            text_frame.margin_bottom = Inches(0.2)
            
            # 제목
            p1 = text_frame.paragraphs[0]
            p1.text = title
            p1.alignment = PP_ALIGN.LEFT
            font1 = p1.font
            font1.size = Pt(14)
            font1.bold = True
            font1.color.rgb = RGBColor(31, 41, 55)
            
            # 설명
            p2 = text_frame.add_paragraph()
            p2.text = description
            p2.alignment = PP_ALIGN.LEFT
            font2 = p2.font
            font2.size = Pt(12)
            font2.color.rgb = RGBColor(107, 114, 128)
            
            return card
            
        except Exception as e:
            print(f"기능 카드 생성 오류: {e}")
            return None
    
    def create_icon_circle(self, slide, x, y, size):
        """아이콘 원형 배경 생성"""
        try:
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x), Inches(y), Inches(size), Inches(size)
            )
            
            # 배경 색상 설정 (blue-100)
            fill = circle.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(219, 234, 254)
            
            return circle
            
        except Exception as e:
            print(f"아이콘 원형 생성 오류: {e}")
            return None
    
    def parse_01_html(self, soup, slide):
        """01.html (메인 페이지) 파싱"""
        # 제목 추가 (text-5xl = 48px)
        title_element = soup.find('h1', class_='title')
        if title_element:
            self.create_text_box(
                slide, title_element.get_text(),
                1, 1, 11, 1.2,
                {'font-size': '48px', 'color': '#2563eb', 'font-weight': '900', 'text-align': 'center'}
            )
        
        # 부제목 추가 (text-3xl = 30px)
        subtitle_element = soup.find('h2', class_='subtitle')
        if subtitle_element:
            self.create_text_box(
                slide, subtitle_element.get_text(),
                1, 2.2, 11, 0.8,
                {'font-size': '30px', 'color': '#1e40af', 'font-weight': '700', 'text-align': 'center'}
            )
        
        # 개발 기간 섹션
        period_label = soup.find('p', string=lambda text: text and '개발 기간' in text)
        if period_label:
            period_value = period_label.find_next_sibling('p')
            if period_value:
                # 라벨
                self.create_text_box(
                    slide, "개발 기간",
                    1, 3.2, 11, 0.4,
                    {'font-size': '20px', 'color': '#6b7280', 'text-align': 'center'}
                )
                # 값
                self.create_text_box(
                    slide, period_value.get_text(),
                    1, 3.6, 11, 0.6,
                    {'font-size': '24px', 'color': '#000000', 'font-weight': '500', 'text-align': 'center'}
                )
        
        # 기술 스택 섹션
        tech_label = soup.find('p', string=lambda text: text and '주요 기술 스택' in text)
        if tech_label:
            # 기술 스택 라벨
            self.create_text_box(
                slide, "주요 기술 스택",
                1, 4.4, 11, 0.4,
                {'font-size': '20px', 'color': '#6b7280', 'text-align': 'center'}
            )
            
            # 기술 스택 컨테이너 찾기
            tech_container = tech_label.find_next('div', class_='flex')
            if tech_container:
                tech_items = tech_container.find_all('div', class_='tech-stack')
                
                # 기술 스택 박스들 (5개를 2행으로 배치)
                start_x = 2.5
                start_y = 5.0
                for i, tech_item in enumerate(tech_items[:5]):
                    tech_text = tech_item.get_text().strip()
                    row = i // 3
                    col = i % 3
                    x = start_x + col * 2.8
                    y = start_y + row * 0.8
                    
                    self.create_tech_stack_box(slide, tech_text, x, y)
        
        # 링크 버튼들
        links_section = soup.find('div', class_='flex justify-center space-x-8')
        if links_section:
            buttons = links_section.find_all('a', class_='link-button')
            
            # GitHub 버튼 (첫 번째)
            if len(buttons) > 0:
                github_button = buttons[0]
                github_text = github_button.find('span').get_text() if github_button.find('span') else 'GitHub'
                self.create_button(
                    slide, github_text,
                    4.5, 6.8, 2.5, 0.6,
                    RGBColor(31, 41, 55), RGBColor(255, 255, 255)  # gray-800, white
                )
            
            # 배포 사이트 버튼 (두 번째)
            if len(buttons) > 1:
                deploy_button = buttons[1]
                deploy_text = deploy_button.find('span').get_text() if deploy_button.find('span') else '배포 사이트'
                self.create_button(
                    slide, deploy_text,
                    7.5, 6.8, 2.5, 0.6,
                    RGBColor(37, 99, 235), RGBColor(255, 255, 255)  # blue-600, white
                )
        
        # 날짜 추가 (우측 하단)
        date_element = soup.find('div', class_='absolute bottom-8 right-8')
        if date_element:
            date_text = date_element.find('p').get_text() if date_element.find('p') else '2025.09.11'
            self.create_text_box(
                slide, date_text,
                10.5, 6.0, 2, 0.4,
                {'font-size': '14px', 'color': '#9ca3af', 'text-align': 'right'}
            )
    
    def parse_02_html(self, soup, slide):
        """02.html (프로젝트 개요) 파싱"""
        # 제목 추가
        title_element = soup.find('h1', class_='section-title')
        if title_element:
            self.create_text_box(
                slide, title_element.get_text(),
                1, 0.5, 11, 1,
                {'font-size': '36px', 'color': '#2563eb', 'font-weight': '800', 'text-align': 'left'}
            )
        
        # 배경 섹션
        background_section = soup.find('h2', string=lambda text: text and '배경' in text)
        if background_section:
            # 아이콘 원형
            self.create_icon_circle(slide, 0.5, 1.8, 0.8)
            
            # 제목
            self.create_text_box(
                slide, "배경",
                1.5, 1.8, 10, 0.6,
                {'font-size': '24px', 'color': '#1f2937', 'font-weight': 'bold', 'text-align': 'left'}
            )
            
            # 내용
            background_content = background_section.find_next('p')
            if background_content:
                self.create_text_box(
                    slide, background_content.get_text(),
                    1.5, 2.4, 10, 1.2,
                    {'font-size': '16px', 'color': '#6b7280', 'text-align': 'left'}
                )
        
        # 목적 섹션
        purpose_section = soup.find('h2', string=lambda text: text and '목적' in text)
        if purpose_section:
            # 아이콘 원형
            self.create_icon_circle(slide, 0.5, 3.8, 0.8)
            
            # 제목
            self.create_text_box(
                slide, "목적",
                1.5, 3.8, 10, 0.6,
                {'font-size': '24px', 'color': '#1f2937', 'font-weight': 'bold', 'text-align': 'left'}
            )
            
            # 내용
            purpose_content = purpose_section.find_next('p')
            if purpose_content:
                self.create_text_box(
                    slide, purpose_content.get_text(),
                    1.5, 4.4, 10, 1.2,
                    {'font-size': '16px', 'color': '#6b7280', 'text-align': 'left'}
                )
        
        # 주요 특징 섹션
        features_section = soup.find('h2', string=lambda text: text and '주요 특징' in text)
        if features_section:
            # 아이콘 원형
            self.create_icon_circle(slide, 0.5, 5.8, 0.8)
            
            # 제목
            self.create_text_box(
                slide, "주요 특징",
                1.5, 5.8, 10, 0.6,
                {'font-size': '24px', 'color': '#1f2937', 'font-weight': 'bold', 'text-align': 'left'}
            )
            
            # 특징 카드들
            feature_cards = soup.find_all('div', class_='feature-card')
            if feature_cards:
                # 2x2 그리드로 배치
                start_x = 1.5
                start_y = 6.4
                card_width = 4.5
                card_height = 1.2
                
                for i, card in enumerate(feature_cards[:4]):
                    row = i // 2
                    col = i % 2
                    x = start_x + col * 5.0
                    y = start_y + row * 1.4
                    
                    title_elem = card.find('h3')
                    desc_elem = card.find('p')
                    
                    if title_elem and desc_elem:
                        title = title_elem.get_text()
                        description = desc_elem.get_text()
                        self.create_feature_card(slide, title, description, x, y, card_width, card_height)
        
        # 푸터
        footer_element = soup.find('div', class_='absolute bottom-8 right-8')
        if footer_element:
            footer_text = footer_element.find('p').get_text() if footer_element.find('p') else '개발 프로젝트: 디지털 창작소 웹사이트'
            self.create_text_box(
                slide, footer_text,
                8, 6.8, 4, 0.4,
                {'font-size': '12px', 'color': '#9ca3af', 'text-align': 'right'}
            )
    
    def parse_html_to_pptx(self):
        """HTML을 파싱하여 PPTX로 변환"""
        try:
            # HTML 파일 읽기
            with open(self.html_file, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # 새 프레젠테이션 생성
            prs = Presentation()
            
            # 슬라이드 크기 설정 (16:9 비율)
            prs.slide_width = Inches(13.33)  # 1920px
            prs.slide_height = Inches(7.5)   # 1080px
            
            # 빈 슬라이드 추가
            slide_layout = prs.slide_layouts[6]  # 빈 레이아웃
            slide = prs.slides.add_slide(slide_layout)
            
            # HTML 파일 타입에 따라 다른 파싱 로직 적용
            if '01.html' in str(self.html_file):
                print("01.html (메인 페이지) 파싱 중...")
                self.parse_01_html(soup, slide)
            elif '02.html' in str(self.html_file):
                print("02.html (프로젝트 개요) 파싱 중...")
                self.parse_02_html(soup, slide)
            else:
                print("알 수 없는 HTML 파일 형식입니다.")
                return False
            
            # PPTX 파일 저장
            prs.save(self.output_path)
            print(f"PPTX 파일 저장 완료: {self.output_path}")
            
            return True
            
        except Exception as e:
            print(f"HTML 파싱 오류: {e}")
            import traceback
            traceback.print_exc()
            return False
    
    def convert(self):
        """전체 변환 프로세스 실행"""
        try:
            print(f"HTML 파일 변환 시작: {self.html_file}")
            
            # 임시 디렉토리 설정
            self.setup_temp_directory()
            
            # HTML을 PPTX로 변환
            success = self.parse_html_to_pptx()
            
            return success
            
        except Exception as e:
            print(f"변환 프로세스 오류: {e}")
            return False
        
        finally:
            # 임시 디렉토리 정리
            self.cleanup_temp_directory()

def main():
    # 설정
    html_file = r"C:\Project\gigabitamin\genspark\dcs_site\html\01.html"
    output_path = r"C:\Project\gigabitamin\genspark\dcs_site\html\01_editable_v3.pptx"
    
    print("HTML to Editable PPTX 변환기 V3 시작")
    print(f"HTML 파일: {html_file}")
    print(f"출력 파일: {output_path}")
    print("-" * 50)
    
    # 변환기 생성 및 실행
    converter = HTMLEditablePPTXConverterV3(html_file, output_path)
    success = converter.convert()
    
    if success:
        print("-" * 50)
        print("변환 완료!")
        print(f"출력 파일: {output_path}")
        print("변환된 PPTX 파일에서 텍스트, 도형, 이미지를 개별적으로 편집할 수 있습니다.")
    else:
        print("-" * 50)
        print("변환 실패!")
        sys.exit(1)

if __name__ == "__main__":
    main()
