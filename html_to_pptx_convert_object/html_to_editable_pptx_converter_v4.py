#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
HTML to Editable PPTX Converter V4
HTML íŒŒì¼ì„ íŒŒì‹±í•˜ì—¬ í…ìŠ¤íŠ¸, ë„í˜•, ì´ë¯¸ì§€ë¥¼ í¸ì§‘ ê°€ëŠ¥í•œ PPTX ê°ì²´ë¡œ ë³€í™˜í•˜ëŠ” ìŠ¤í¬ë¦½íŠ¸
- ë ˆì´ì•„ì›ƒ ìµœì í™” (ì˜ë¦¼ í˜„ìƒ í•´ê²°)
- ì•„ì´ì½˜ í…ìŠ¤íŠ¸ ì§€ì› ì¶”ê°€
"""

import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup
import re
import requests
from urllib.parse import urljoin, urlparse
import tempfile
import shutil

class HTMLEditablePPTXConverterV4:
    def __init__(self, html_file, output_path):
        self.html_file = Path(html_file)
        self.output_path = Path(output_path)
        self.temp_dir = None
        
    def setup_temp_directory(self):
        """ì„ì‹œ ë””ë ‰í† ë¦¬ ì„¤ì •"""
        self.temp_dir = Path(tempfile.mkdtemp())
        print(f"ì„ì‹œ ë””ë ‰í† ë¦¬ ìƒì„±: {self.temp_dir}")
        
    def cleanup_temp_directory(self):
        """ì„ì‹œ ë””ë ‰í† ë¦¬ ì •ë¦¬"""
        if self.temp_dir and self.temp_dir.exists():
            shutil.rmtree(self.temp_dir)
            print("ì„ì‹œ ë””ë ‰í† ë¦¬ ì •ë¦¬ ì™„ë£Œ")
    
    def parse_css_color(self, color_str):
        """CSS ìƒ‰ìƒ ë¬¸ìì—´ì„ RGBColorë¡œ ë³€í™˜"""
        if not color_str:
            return None
            
        # #RRGGBB í˜•ì‹
        if color_str.startswith('#'):
            hex_color = color_str[1:]
            if len(hex_color) == 6:
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
                return RGBColor(r, g, b)
        
        # rgb(r, g, b) í˜•ì‹
        rgb_match = re.match(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', color_str)
        if rgb_match:
            r, g, b = map(int, rgb_match.groups())
            return RGBColor(r, g, b)
        
        # ìƒ‰ìƒ ì´ë¦„ ë§¤í•‘
        color_map = {
            'blue': RGBColor(37, 99, 235),
            'red': RGBColor(239, 68, 68),
            'green': RGBColor(34, 197, 94),
            'yellow': RGBColor(234, 179, 8),
            'purple': RGBColor(147, 51, 234),
            'gray': RGBColor(107, 114, 128),
            'black': RGBColor(0, 0, 0),
            'white': RGBColor(255, 255, 255)
        }
        
        return color_map.get(color_str.lower())
    
    def parse_font_size(self, font_size_str):
        """CSS í°íŠ¸ í¬ê¸°ë¥¼ Ptë¡œ ë³€í™˜"""
        if not font_size_str:
            return 12
        
        # px ë‹¨ìœ„ ì œê±°
        if font_size_str.endswith('px'):
            return int(float(font_size_str[:-2]))
        elif font_size_str.endswith('rem'):
            return int(float(font_size_str[:-3]) * 16)  # 1rem = 16px ê°€ì •
        elif font_size_str.endswith('em'):
            return int(float(font_size_str[:-2]) * 16)  # 1em = 16px ê°€ì •
        else:
            try:
                return int(float(font_size_str))
            except:
                return 12
    
    def get_text_alignment(self, text_align):
        """CSS text-alignì„ PPTX ì •ë ¬ë¡œ ë³€í™˜"""
        align_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        return align_map.get(text_align, PP_ALIGN.LEFT)
    
    def get_icon_text(self, icon_class):
        """FontAwesome ì•„ì´ì½˜ í´ë˜ìŠ¤ë¥¼ í…ìŠ¤íŠ¸ë¡œ ë³€í™˜"""
        # fas, fab, far ì ‘ë‘ì‚¬ ì œê±°
        clean_class = icon_class.replace('fas ', '').replace('fab ', '').replace('far ', '')
        
        icon_map = {
            'fa-react': 'âš›',
            'fa-js': 'JS',
            'fa-css3': 'CSS',
            'fa-database': 'ğŸ—„',
            'fa-server': 'ğŸ–¥',
            'fa-github': 'ğŸ™',
            'fa-globe': 'ğŸŒ',
            'fa-history': 'â˜‘â˜‘',  # ì‹¤ì œ ë Œë”ë§ëœ ëª¨ì–‘ìœ¼ë¡œ ë³€ê²½
            'fa-bullseye': 'â˜‘â˜‘',  # ì‹¤ì œ ë Œë”ë§ëœ ëª¨ì–‘ìœ¼ë¡œ ë³€ê²½
            'fa-star': 'â­',
            'fa-users': 'ğŸ‘¥',
            'fa-graduation-cap': 'ğŸ“',
            'fa-project-diagram': 'ğŸ“Š',
            'fa-mobile-alt': 'ğŸ“±'
        }
        return icon_map.get(clean_class, 'â—')
    
    def create_text_box(self, slide, text, x, y, width, height, styles=None):
        """í…ìŠ¤íŠ¸ ë°•ìŠ¤ ìƒì„±"""
        try:
            # í…ìŠ¤íŠ¸ ë°•ìŠ¤ ì¶”ê°€
            textbox = slide.shapes.add_textbox(
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            
            # í…ìŠ¤íŠ¸ í”„ë ˆì„ ì„¤ì •
            text_frame = textbox.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.05)
            text_frame.margin_bottom = Inches(0.05)
            
            # ë‹¨ë½ ìƒì„±
            p = text_frame.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.LEFT
            
            # ìŠ¤íƒ€ì¼ ì ìš©
            if styles:
                font = p.font
                
                # í°íŠ¸ í¬ê¸°
                if 'font-size' in styles:
                    font.size = Pt(self.parse_font_size(styles['font-size']))
                
                # í°íŠ¸ ìƒ‰ìƒ
                if 'color' in styles:
                    color = self.parse_css_color(styles['color'])
                    if color:
                        font.color.rgb = color
                
                # í°íŠ¸ êµµê¸°
                if 'font-weight' in styles:
                    if styles['font-weight'] in ['bold', '700', '800', '900']:
                        font.bold = True
                
                # ì •ë ¬
                if 'text-align' in styles:
                    p.alignment = self.get_text_alignment(styles['text-align'])
            
            return textbox
            
        except Exception as e:
            print(f"í…ìŠ¤íŠ¸ ë°•ìŠ¤ ìƒì„± ì˜¤ë¥˜: {e}")
            return None
    
    def create_tech_stack_box(self, slide, tech_text, icon_text, x, y):
        """ê¸°ìˆ  ìŠ¤íƒ ë°•ìŠ¤ ìƒì„±"""
        try:
            # ë°°ê²½ ë°•ìŠ¤ ìƒì„±
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(2.2), Inches(0.6)
            )
            
            # ë°°ê²½ ìƒ‰ìƒ ì„¤ì • (blue-50)
            fill = box.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(239, 246, 255)
            
            # í…Œë‘ë¦¬ ì„¤ì •
            line = box.line
            line.color.rgb = RGBColor(219, 234, 254)
            
            # í…ìŠ¤íŠ¸ ì¶”ê°€
            text_frame = box.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            p = text_frame.paragraphs[0]
            p.text = f"{icon_text} {tech_text}"
            p.alignment = PP_ALIGN.CENTER
            
            # í°íŠ¸ ì„¤ì •
            font = p.font
            font.size = Pt(12)
            font.bold = True
            font.color.rgb = RGBColor(30, 64, 175)  # blue-800
            
            return box
            
        except Exception as e:
            print(f"ê¸°ìˆ  ìŠ¤íƒ ë°•ìŠ¤ ìƒì„± ì˜¤ë¥˜: {e}")
            return None
    
    def create_button(self, slide, text, icon_text, x, y, width, height, bg_color, text_color):
        """ë²„íŠ¼ ìƒì„±"""
        try:
            # ë²„íŠ¼ ë°•ìŠ¤ ìƒì„±
            button = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            
            # ë°°ê²½ ìƒ‰ìƒ ì„¤ì •
            fill = button.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
            
            # í…ìŠ¤íŠ¸ ì¶”ê°€
            text_frame = button.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            p = text_frame.paragraphs[0]
            p.text = f"{icon_text} {text}"
            p.alignment = PP_ALIGN.CENTER
            
            # í°íŠ¸ ì„¤ì •
            font = p.font
            font.size = Pt(14)
            font.bold = True
            font.color.rgb = text_color
            
            return button
            
        except Exception as e:
            print(f"ë²„íŠ¼ ìƒì„± ì˜¤ë¥˜: {e}")
            return None
    
    def create_feature_card(self, slide, title, description, icon_text, x, y, width, height):
        """ê¸°ëŠ¥ ì¹´ë“œ ìƒì„±"""
        try:
            # ì¹´ë“œ ë°°ê²½ ìƒì„±
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            
            # ë°°ê²½ ìƒ‰ìƒ ì„¤ì • (gray-50)
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(249, 250, 251)
            
            # í…Œë‘ë¦¬ ì„¤ì •
            line = card.line
            line.color.rgb = RGBColor(229, 231, 235)
            
            # í…ìŠ¤íŠ¸ ì¶”ê°€
            text_frame = card.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)
            text_frame.margin_top = Inches(0.2)
            text_frame.margin_bottom = Inches(0.2)
            
            # ì œëª©
            p1 = text_frame.paragraphs[0]
            p1.text = f"{icon_text} {title}"
            p1.alignment = PP_ALIGN.LEFT
            font1 = p1.font
            font1.size = Pt(14)
            font1.bold = True
            font1.color.rgb = RGBColor(31, 41, 55)
            
            # ì„¤ëª…
            p2 = text_frame.add_paragraph()
            p2.text = description
            p2.alignment = PP_ALIGN.LEFT
            font2 = p2.font
            font2.size = Pt(12)
            font2.color.rgb = RGBColor(107, 114, 128)
            
            return card
            
        except Exception as e:
            print(f"ê¸°ëŠ¥ ì¹´ë“œ ìƒì„± ì˜¤ë¥˜: {e}")
            return None
    
    def create_icon_circle(self, slide, icon_text, x, y, size):
        """ì•„ì´ì½˜ ì›í˜• ë°°ê²½ ìƒì„±"""
        try:
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x), Inches(y), Inches(size), Inches(size)
            )
            
            # ë°°ê²½ ìƒ‰ìƒ ì„¤ì • (blue-100)
            fill = circle.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(219, 234, 254)
            
            # ì•„ì´ì½˜ í…ìŠ¤íŠ¸ ì¶”ê°€
            text_frame = circle.text_frame
            text_frame.clear()
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            p = text_frame.paragraphs[0]
            p.text = icon_text
            p.alignment = PP_ALIGN.CENTER
            
            font = p.font
            font.size = Pt(20)
            font.color.rgb = RGBColor(37, 99, 235)
            
            return circle
            
        except Exception as e:
            print(f"ì•„ì´ì½˜ ì›í˜• ìƒì„± ì˜¤ë¥˜: {e}")
            return None
    
    def parse_01_html(self, soup, slide):
        """01.html (ë©”ì¸ í˜ì´ì§€) íŒŒì‹±"""
        # ì œëª© ì¶”ê°€ (text-5xl = 48px)
        title_element = soup.find('h1', class_='title')
        if title_element:
            self.create_text_box(
                slide, title_element.get_text(),
                1, 0.8, 11, 1,
                {'font-size': '48px', 'color': '#2563eb', 'font-weight': '900', 'text-align': 'center'}
            )
        
        # ë¶€ì œëª© ì¶”ê°€ (text-3xl = 30px)
        subtitle_element = soup.find('h2', class_='subtitle')
        if subtitle_element:
            self.create_text_box(
                slide, subtitle_element.get_text(),
                1, 1.8, 11, 0.6,
                {'font-size': '30px', 'color': '#1e40af', 'font-weight': '700', 'text-align': 'center'}
            )
        
        # ê°œë°œ ê¸°ê°„ ì„¹ì…˜
        period_label = soup.find('p', string=lambda text: text and 'ê°œë°œ ê¸°ê°„' in text)
        if period_label:
            period_value = period_label.find_next_sibling('p')
            if period_value:
                # ë¼ë²¨
                self.create_text_box(
                    slide, "ê°œë°œ ê¸°ê°„",
                    1, 2.6, 11, 0.3,
                    {'font-size': '20px', 'color': '#6b7280', 'text-align': 'center'}
                )
                # ê°’
                self.create_text_box(
                    slide, period_value.get_text(),
                    1, 2.9, 11, 0.5,
                    {'font-size': '24px', 'color': '#000000', 'font-weight': '500', 'text-align': 'center'}
                )
        
        # ê¸°ìˆ  ìŠ¤íƒ ì„¹ì…˜
        tech_label = soup.find('p', string=lambda text: text and 'ì£¼ìš” ê¸°ìˆ  ìŠ¤íƒ' in text)
        if tech_label:
            # ê¸°ìˆ  ìŠ¤íƒ ë¼ë²¨
            self.create_text_box(
                slide, "ì£¼ìš” ê¸°ìˆ  ìŠ¤íƒ",
                1, 3.6, 11, 0.3,
                {'font-size': '20px', 'color': '#6b7280', 'text-align': 'center'}
            )
            
            # ê¸°ìˆ  ìŠ¤íƒ ì»¨í…Œì´ë„ˆ ì°¾ê¸°
            tech_container = tech_label.find_next('div', class_='flex')
            if tech_container:
                tech_items = tech_container.find_all('div', class_='tech-stack')
                
                # ê¸°ìˆ  ìŠ¤íƒ ë°•ìŠ¤ë“¤ (5ê°œë¥¼ 2í–‰ìœ¼ë¡œ ë°°ì¹˜)
                start_x = 2.5
                start_y = 4.1
                for i, tech_item in enumerate(tech_items[:5]):
                    tech_text = tech_item.get_text().strip()
                    # ì•„ì´ì½˜ ì¶”ì¶œ
                    icon_elem = tech_item.find('i')
                    icon_text = 'â—'
                    if icon_elem:
                        icon_classes = icon_elem.get('class', [])
                        for cls in icon_classes:
                            if cls.startswith('fa-'):
                                icon_text = self.get_icon_text(cls)
                                break
                    
                    row = i // 3
                    col = i % 3
                    x = start_x + col * 2.8
                    y = start_y + row * 0.8
                    
                    self.create_tech_stack_box(slide, tech_text, icon_text, x, y)
        
        # ë§í¬ ë²„íŠ¼ë“¤
        links_section = soup.find('div', class_='flex justify-center space-x-8')
        if links_section:
            buttons = links_section.find_all('a', class_='link-button')
            
            # GitHub ë²„íŠ¼ (ì²« ë²ˆì§¸)
            if len(buttons) > 0:
                github_button = buttons[0]
                github_text = github_button.find('span').get_text() if github_button.find('span') else 'GitHub'
                github_icon = github_button.find('i', class_='fab fa-github')
                github_icon_text = self.get_icon_text('fa-github') if github_icon else 'ğŸ™'
                self.create_button(
                    slide, github_text, github_icon_text,
                    4.5, 6.2, 2.5, 0.6,
                    RGBColor(31, 41, 55), RGBColor(255, 255, 255)  # gray-800, white
                )
            
            # ë°°í¬ ì‚¬ì´íŠ¸ ë²„íŠ¼ (ë‘ ë²ˆì§¸)
            if len(buttons) > 1:
                deploy_button = buttons[1]
                deploy_text = deploy_button.find('span').get_text() if deploy_button.find('span') else 'ë°°í¬ ì‚¬ì´íŠ¸'
                deploy_icon = deploy_button.find('i', class_='fas fa-globe')
                deploy_icon_text = self.get_icon_text('fa-globe') if deploy_icon else 'ğŸŒ'
                self.create_button(
                    slide, deploy_text, deploy_icon_text,
                    7.5, 6.2, 2.5, 0.6,
                    RGBColor(37, 99, 235), RGBColor(255, 255, 255)  # blue-600, white
                )
        
        # ë‚ ì§œ ì¶”ê°€ (ìš°ì¸¡ í•˜ë‹¨)
        date_element = soup.find('div', class_='absolute bottom-8 right-8')
        if date_element:
            date_text = date_element.find('p').get_text() if date_element.find('p') else '2025.09.11'
            self.create_text_box(
                slide, date_text,
                10.5, 5.5, 2, 0.3,
                {'font-size': '14px', 'color': '#9ca3af', 'text-align': 'right'}
            )
    
    def parse_02_html(self, soup, slide):
        """02.html (í”„ë¡œì íŠ¸ ê°œìš”) íŒŒì‹±"""
        # ì œëª© ì¶”ê°€
        title_element = soup.find('h1', class_='section-title')
        if title_element:
            self.create_text_box(
                slide, title_element.get_text(),
                1, 0.3, 11, 0.8,
                {'font-size': '36px', 'color': '#2563eb', 'font-weight': '800', 'text-align': 'left'}
            )
        
        # ë°°ê²½ ì„¹ì…˜
        background_section = soup.find('h2', string=lambda text: text and 'ë°°ê²½' in text)
        if background_section:
            # ì•„ì´ì½˜ ì°¾ê¸° - ë¶€ëª¨ divì—ì„œ ì°¾ê¸°
            parent_div = background_section.find_parent('div', class_='flex')
            icon_elem = None
            if parent_div:
                icon_elem = parent_div.find('i')
            
            icon_text = 'ğŸ“š'  # ê¸°ë³¸ê°’
            if icon_elem:
                icon_classes = icon_elem.get('class', [])
                for cls in icon_classes:
                    if cls.startswith('fa-'):
                        icon_text = self.get_icon_text(cls)
                        break
            
            # ì•„ì´ì½˜ ì›í˜•
            self.create_icon_circle(slide, icon_text, 0.5, 1.2, 0.6)
            
            # ì œëª©
            self.create_text_box(
                slide, "ë°°ê²½",
                1.3, 1.2, 10, 0.5,
                {'font-size': '24px', 'color': '#1f2937', 'font-weight': 'bold', 'text-align': 'left'}
            )
            
            # ë‚´ìš©
            background_content = background_section.find_next('p')
            if background_content:
                self.create_text_box(
                    slide, background_content.get_text(),
                    1.3, 1.7, 10, 1,
                    {'font-size': '16px', 'color': '#6b7280', 'text-align': 'left'}
                )
        
        # ëª©ì  ì„¹ì…˜
        purpose_section = soup.find('h2', string=lambda text: text and 'ëª©ì ' in text)
        if purpose_section:
            # ì•„ì´ì½˜ ì°¾ê¸° - ë¶€ëª¨ divì—ì„œ ì°¾ê¸°
            parent_div = purpose_section.find_parent('div', class_='flex')
            icon_elem = None
            if parent_div:
                icon_elem = parent_div.find('i')
            
            icon_text = 'ğŸ¯'  # ê¸°ë³¸ê°’
            if icon_elem:
                icon_classes = icon_elem.get('class', [])
                for cls in icon_classes:
                    if cls.startswith('fa-'):
                        icon_text = self.get_icon_text(cls)
                        break
            
            # ì•„ì´ì½˜ ì›í˜•
            self.create_icon_circle(slide, icon_text, 0.5, 2.9, 0.6)
            
            # ì œëª©
            self.create_text_box(
                slide, "ëª©ì ",
                1.3, 2.9, 10, 0.5,
                {'font-size': '24px', 'color': '#1f2937', 'font-weight': 'bold', 'text-align': 'left'}
            )
            
            # ë‚´ìš©
            purpose_content = purpose_section.find_next('p')
            if purpose_content:
                self.create_text_box(
                    slide, purpose_content.get_text(),
                    1.3, 3.4, 10, 1,
                    {'font-size': '16px', 'color': '#6b7280', 'text-align': 'left'}
                )
        
        # ì£¼ìš” íŠ¹ì§• ì„¹ì…˜
        features_section = soup.find('h2', string=lambda text: text and 'ì£¼ìš” íŠ¹ì§•' in text)
        if features_section:
            # ì•„ì´ì½˜ ì°¾ê¸° - ë¶€ëª¨ divì—ì„œ ì°¾ê¸°
            parent_div = features_section.find_parent('div', class_='flex')
            icon_elem = None
            if parent_div:
                icon_elem = parent_div.find('i')
            
            icon_text = 'â­'  # ê¸°ë³¸ê°’
            if icon_elem:
                icon_classes = icon_elem.get('class', [])
                for cls in icon_classes:
                    if cls.startswith('fa-'):
                        icon_text = self.get_icon_text(cls)
                        break
            
            # ì•„ì´ì½˜ ì›í˜•
            self.create_icon_circle(slide, icon_text, 0.5, 4.6, 0.6)
            
            # ì œëª©
            self.create_text_box(
                slide, "ì£¼ìš” íŠ¹ì§•",
                1.3, 4.6, 10, 0.5,
                {'font-size': '24px', 'color': '#1f2937', 'font-weight': 'bold', 'text-align': 'left'}
            )
            
            # íŠ¹ì§• ì¹´ë“œë“¤
            feature_cards = soup.find_all('div', class_='feature-card')
            if feature_cards:
                # 2x2 ê·¸ë¦¬ë“œë¡œ ë°°ì¹˜
                start_x = 1.3
                start_y = 5.2
                card_width = 4.5
                card_height = 1.0
                
                for i, card in enumerate(feature_cards[:4]):
                    row = i // 2
                    col = i % 2
                    x = start_x + col * 5.0
                    y = start_y + row * 1.2
                    
                    title_elem = card.find('h3')
                    desc_elem = card.find('p')
                    icon_elem = card.find('i')
                    
                    if title_elem and desc_elem:
                        title = title_elem.get_text()
                        description = desc_elem.get_text()
                        
                        # ì•„ì´ì½˜ ì¶”ì¶œ
                        icon_text = 'â—'
                        if icon_elem:
                            icon_classes = icon_elem.get('class', [])
                            for cls in icon_classes:
                                if cls.startswith('fa-'):
                                    icon_text = self.get_icon_text(cls)
                                    break
                        
                        self.create_feature_card(slide, title, description, icon_text, x, y, card_width, card_height)
        
        # í‘¸í„°
        footer_element = soup.find('div', class_='absolute bottom-8 right-8')
        if footer_element:
            footer_text = footer_element.find('p').get_text() if footer_element.find('p') else 'ê°œë°œ í”„ë¡œì íŠ¸: ë””ì§€í„¸ ì°½ì‘ì†Œ ì›¹ì‚¬ì´íŠ¸'
            self.create_text_box(
                slide, footer_text,
                8, 6.8, 4, 0.3,
                {'font-size': '12px', 'color': '#9ca3af', 'text-align': 'right'}
            )
    
    def parse_html_to_pptx(self):
        """HTMLì„ íŒŒì‹±í•˜ì—¬ PPTXë¡œ ë³€í™˜"""
        try:
            # HTML íŒŒì¼ ì½ê¸°
            with open(self.html_file, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # ìƒˆ í”„ë ˆì  í…Œì´ì…˜ ìƒì„±
            prs = Presentation()
            
            # ìŠ¬ë¼ì´ë“œ í¬ê¸° ì„¤ì • (16:9 ë¹„ìœ¨)
            prs.slide_width = Inches(13.33)  # 1920px
            prs.slide_height = Inches(7.5)   # 1080px
            
            # ë¹ˆ ìŠ¬ë¼ì´ë“œ ì¶”ê°€
            slide_layout = prs.slide_layouts[6]  # ë¹ˆ ë ˆì´ì•„ì›ƒ
            slide = prs.slides.add_slide(slide_layout)
            
            # HTML íŒŒì¼ íƒ€ì…ì— ë”°ë¼ ë‹¤ë¥¸ íŒŒì‹± ë¡œì§ ì ìš©
            if '01.html' in str(self.html_file):
                print("01.html (ë©”ì¸ í˜ì´ì§€) íŒŒì‹± ì¤‘...")
                self.parse_01_html(soup, slide)
            elif '02.html' in str(self.html_file):
                print("02.html (í”„ë¡œì íŠ¸ ê°œìš”) íŒŒì‹± ì¤‘...")
                self.parse_02_html(soup, slide)
            else:
                print("ì•Œ ìˆ˜ ì—†ëŠ” HTML íŒŒì¼ í˜•ì‹ì…ë‹ˆë‹¤.")
                return False
            
            # PPTX íŒŒì¼ ì €ì¥
            prs.save(self.output_path)
            print(f"PPTX íŒŒì¼ ì €ì¥ ì™„ë£Œ: {self.output_path}")
            
            return True
            
        except Exception as e:
            print(f"HTML íŒŒì‹± ì˜¤ë¥˜: {e}")
            import traceback
            traceback.print_exc()
            return False
    
    def convert(self):
        """ì „ì²´ ë³€í™˜ í”„ë¡œì„¸ìŠ¤ ì‹¤í–‰"""
        try:
            print(f"HTML íŒŒì¼ ë³€í™˜ ì‹œì‘: {self.html_file}")
            
            # ì„ì‹œ ë””ë ‰í† ë¦¬ ì„¤ì •
            self.setup_temp_directory()
            
            # HTMLì„ PPTXë¡œ ë³€í™˜
            success = self.parse_html_to_pptx()
            
            return success
            
        except Exception as e:
            print(f"ë³€í™˜ í”„ë¡œì„¸ìŠ¤ ì˜¤ë¥˜: {e}")
            return False
        
        finally:
            # ì„ì‹œ ë””ë ‰í† ë¦¬ ì •ë¦¬
            self.cleanup_temp_directory()

def main():
    # ì„¤ì •
    html_file = r"C:\Project\gigabitamin\genspark\dcs_site\html\01.html"
    output_path = r"C:\Project\gigabitamin\genspark\dcs_site\html\01_editable_v4.pptx"
    
    print("HTML to Editable PPTX ë³€í™˜ê¸° V4 ì‹œì‘")
    print(f"HTML íŒŒì¼: {html_file}")
    print(f"ì¶œë ¥ íŒŒì¼: {output_path}")
    print("-" * 50)
    
    # ë³€í™˜ê¸° ìƒì„± ë° ì‹¤í–‰
    converter = HTMLEditablePPTXConverterV4(html_file, output_path)
    success = converter.convert()
    
    if success:
        print("-" * 50)
        print("ë³€í™˜ ì™„ë£Œ!")
        print(f"ì¶œë ¥ íŒŒì¼: {output_path}")
        print("ë³€í™˜ëœ PPTX íŒŒì¼ì—ì„œ í…ìŠ¤íŠ¸, ë„í˜•, ì´ë¯¸ì§€ë¥¼ ê°œë³„ì ìœ¼ë¡œ í¸ì§‘í•  ìˆ˜ ìˆìŠµë‹ˆë‹¤.")
    else:
        print("-" * 50)
        print("ë³€í™˜ ì‹¤íŒ¨!")
        sys.exit(1)

if __name__ == "__main__":
    main()
