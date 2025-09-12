#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
HTML to Editable PPTX Converter V6
HTML 파일을 파싱하여 텍스트, 도형, 이미지를 편집 가능한 PPTX 객체로 변환하는 스크립트
- FontAwesome 아이콘을 HTML 스크린샷으로 캡처하여 삽입
"""

import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup
import re
import requests
from urllib.parse import urljoin, urlparse
import base64
import tempfile
import shutil
from html2image import Html2Image
from PIL import Image

class HTMLEditablePPTXConverterV6:
    def __init__(self, html_file, output_path):
        self.html_file = Path(html_file)
        self.output_path = Path(output_path)
        self.temp_dir = None
        self.hti = Html2Image()
        
    def setup_temp_directory(self):
        """임시 디렉토리 설정"""
        self.temp_dir = Path(tempfile.mkdtemp())
        print(f"임시 디렉토리 생성: {self.temp_dir}")
        
    def cleanup_temp_directory(self):
        """임시 디렉토리 정리"""
        if self.temp_dir and self.temp_dir.exists():
            shutil.rmtree(self.temp_dir)
            print("임시 디렉토리 정리 완료")
    
    def parse_css_color(self, color_str):
        """CSS 색상 문자열을 RGBColor로 변환"""
        if not color_str:
            return None
            
        # #RRGGBB 형식
        if color_str.startswith('#'):
            hex_color = color_str[1:]
            if len(hex_color) == 6:
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
                return RGBColor(r, g, b)
        
        # rgb(r, g, b) 형식
        rgb_match = re.match(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', color_str)
        if rgb_match:
            r, g, b = map(int, rgb_match.groups())
            return RGBColor(r, g, b)
        
        # 색상 이름 매핑
        color_map = {
            'blue': RGBColor(37, 99, 235),
            'red': RGBColor(239, 68, 68),
            'green': RGBColor(34, 197, 94),
            'yellow': RGBColor(234, 179, 8),
            'purple': RGBColor(147, 51, 234),
            'gray': RGBColor(107, 114, 128),
            'black': RGBColor(0, 0, 0),
            'white': RGBColor(255, 255, 255)
        }
        
        return color_map.get(color_str.lower())
    
    def parse_font_size(self, font_size_str):
        """CSS 폰트 크기를 Pt로 변환"""
        if not font_size_str:
            return 12
        
        # px 단위 제거
        if font_size_str.endswith('px'):
            return int(float(font_size_str[:-2]))
        elif font_size_str.endswith('rem'):
            return int(float(font_size_str[:-3]) * 16)  # 1rem = 16px 가정
        elif font_size_str.endswith('em'):
            return int(float(font_size_str[:-2]) * 16)  # 1em = 16px 가정
        else:
            try:
                return int(float(font_size_str))
            except:
                return 12
    
    def get_text_alignment(self, text_align):
        """CSS text-align을 PPTX 정렬로 변환"""
        align_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        return align_map.get(text_align, PP_ALIGN.LEFT)
    
    def create_icon_html(self, icon_class, color='#2563eb', size=64):
        """FontAwesome 아이콘을 위한 HTML 생성"""
        return f"""
        <!DOCTYPE html>
        <html>
        <head>
            <link href="https://cdn.jsdelivr.net/npm/@fortawesome/fontawesome-free@6.4.0/css/all.min.css" rel="stylesheet"/>
            <style>
                body {{
                    margin: 0;
                    padding: 0;
                    background: transparent;
                    display: flex;
                    justify-content: center;
                    align-items: center;
                    width: {size}px;
                    height: {size}px;
                }}
                .icon {{
                    font-size: {size * 0.6}px;
                    color: {color};
                }}
            </style>
        </head>
        <body>
            <i class="{icon_class} icon"></i>
        </body>
        </html>
        """
    
    def download_fontawesome_svg(self, icon_class, color='#2563eb'):
        """FontAwesome 아이콘 SVG 다운로드"""
        try:
            clean_class = icon_class.replace('fas ', '').replace('fab ', '').replace('far ', '').replace('fa-', '')
            
            # 다양한 경로 시도
            paths = [
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/6.x/svgs/solid/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/6.x/svgs/brands/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/6.x/svgs/regular/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/5.x/svgs/solid/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/5.x/svgs/brands/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/5.x/svgs/regular/{clean_class}.svg",
            ]
            
            for path in paths:
                try:
                    print(f"아이콘 다운로드 시도: {icon_class} -> {path}")
                    response = requests.get(path, timeout=10)
                    if response.status_code == 200:
                        svg_content = response.text
                        
                        # 색상 적용
                        svg_content = svg_content.replace('fill="currentColor"', f'fill="{color}"')
                        svg_content = svg_content.replace('fill="#000"', f'fill="{color}"')
                        svg_content = svg_content.replace('fill="black"', f'fill="{color}"')
                        
                        # SVG 파일 저장
                        svg_file = self.temp_dir / f"{clean_class}.svg"
                        with open(svg_file, 'w', encoding='utf-8') as f:
                            f.write(svg_content)
                        
                        print(f"아이콘 다운로드 성공: {svg_file}")
                        return svg_file
                        
                except Exception as e:
                    print(f"경로 실패: {e}")
                    continue
            
            print(f"아이콘 다운로드 실패: {icon_class}")
            return None
            
        except Exception as e:
            print(f"아이콘 다운로드 오류 ({icon_class}): {e}")
            return None
    
    def svg_to_png_with_html2image(self, svg_file, size=64):
        """SVG를 HTML2Image로 PNG 변환"""
        try:
            # SVG 파일 읽기
            with open(svg_file, 'r', encoding='utf-8') as f:
                svg_content = f.read()
            
            # HTML 생성
            html_content = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <style>
                    body {{
                        margin: 0;
                        padding: 0;
                        width: {size}px;
                        height: {size}px;
                        display: flex;
                        align-items: center;
                        justify-content: center;
                        background: transparent;
                    }}
                    svg {{
                        width: {size}px;
                        height: {size}px;
                    }}
                </style>
            </head>
            <body>
                {svg_content}
            </body>
            </html>
            """
            
            # 임시 HTML 파일 생성
            temp_html = self.temp_dir / f"{svg_file.stem}_temp.html"
            with open(temp_html, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            # PNG 파일 경로
            png_file = self.temp_dir / f"{svg_file.stem}.png"
            
            # HTML2Image로 변환
            self.hti.screenshot(
                html_file=str(temp_html),
                save_as=f"{svg_file.stem}.png",
                size=(size, size)
            )
            
            # 생성된 PNG 파일을 임시 디렉토리로 이동
            generated_png = Path(f"{svg_file.stem}.png")
            if generated_png.exists():
                shutil.move(str(generated_png), str(png_file))
                print(f"SVG to PNG 변환 성공: {png_file}")
                return png_file
            else:
                print(f"PNG 파일 생성 실패: {generated_png}")
                return None
                
        except Exception as e:
            print(f"SVG to PNG 변환 실패: {e}")
            return None
    
    def create_icon_text_box(self, slide, icon_text, x, y, size=0.5):
        """아이콘 텍스트 박스 생성"""
        try:
            # 원형 배경 생성
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x, y, Inches(size), Inches(size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(37, 99, 235)  # 파란색
            
            # 텍스트 프레임 추가
            text_frame = circle.text_frame
            text_frame.text = icon_text
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # 텍스트 스타일 설정
            p = text_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)  # 흰색
            p.font.bold = True
            
            return circle
            
        except Exception as e:
            print(f"아이콘 텍스트 박스 생성 실패: {e}")
            return None
    
    def create_text_box(self, slide, text, x, y, width, height, styles=None):
        """텍스트 박스 생성"""
        try:
            # 텍스트 박스 추가
            textbox = slide.shapes.add_textbox(
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            
            # 텍스트 프레임 설정
            text_frame = textbox.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.05)
            text_frame.margin_bottom = Inches(0.05)
            
            # 단락 생성
            p = text_frame.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.LEFT
            
            # 스타일 적용
            if styles:
                font = p.font
                
                # 폰트 크기
                if 'font-size' in styles:
                    font.size = Pt(self.parse_font_size(styles['font-size']))
                
                # 폰트 색상
                if 'color' in styles:
                    color = self.parse_css_color(styles['color'])
                    if color:
                        font.color.rgb = color
                
                # 폰트 굵기
                if 'font-weight' in styles:
                    if styles['font-weight'] in ['bold', '700', '800', '900']:
                        font.bold = True
                
                # 정렬
                if 'text-align' in styles:
                    p.alignment = self.get_text_alignment(styles['text-align'])
            
            return textbox
            
        except Exception as e:
            print(f"텍스트 박스 생성 오류: {e}")
            return None
    
    def create_tech_stack_box(self, slide, tech_text, icon_class, x, y):
        """기술 스택 박스 생성 (이미지 아이콘 버전)"""
        try:
            # 배경 박스 생성
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(2.2), Inches(0.6)
            )
            
            # 배경 색상 설정 (blue-50)
            fill = box.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(239, 246, 255)
            
            # 테두리 설정
            line = box.line
            line.color.rgb = RGBColor(219, 234, 254)
            
            # 텍스트 추가
            text_frame = box.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            p = text_frame.paragraphs[0]
            p.text = tech_text
            p.alignment = PP_ALIGN.CENTER
            
            # 폰트 설정
            font = p.font
            font.size = Pt(12)
            font.bold = True
            font.color.rgb = RGBColor(30, 64, 175)  # blue-800
            
            # 아이콘 이미지 추가
            if icon_class:
                svg_file = self.download_fontawesome_svg(icon_class, '#1e40af')
                if svg_file:
                    png_file = self.svg_to_png_with_html2image(svg_file, 24)
                    if png_file and png_file.exists():
                        # 아이콘 이미지 추가
                        icon_left = Inches(x + 0.1)
                        icon_top = Inches(y + 0.1)
                        icon_width = Inches(0.3)
                        icon_height = Inches(0.3)
                        
                        slide.shapes.add_picture(str(png_file), icon_left, icon_top, icon_width, icon_height)
                        print(f"기술 스택 아이콘 추가 성공: {icon_class}")
                    else:
                        print(f"기술 스택 PNG 변환 실패: {icon_class}")
                else:
                    print(f"기술 스택 SVG 다운로드 실패: {icon_class}")
            
            return box
            
        except Exception as e:
            print(f"기술 스택 박스 생성 오류: {e}")
            return None
    
    def create_button(self, slide, text, icon_class, x, y, width, height, bg_color, text_color):
        """버튼 생성 (이미지 아이콘 버전)"""
        try:
            # 버튼 박스 생성
            button = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            
            # 배경 색상 설정
            fill = button.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
            
            # 텍스트 추가
            text_frame = button.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            p = text_frame.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.CENTER
            
            # 폰트 설정
            font = p.font
            font.size = Pt(14)
            font.bold = True
            font.color.rgb = text_color
            
            # 아이콘 이미지 추가
            if icon_class:
                svg_file = self.download_fontawesome_svg(icon_class, '#ffffff')
                if svg_file:
                    png_file = self.svg_to_png_with_html2image(svg_file, 24)
                    if png_file and png_file.exists():
                        # 아이콘 이미지 추가
                        icon_left = Inches(x + 0.1)
                        icon_top = Inches(y + 0.1)
                        icon_width = Inches(0.3)
                        icon_height = Inches(0.3)
                        
                        slide.shapes.add_picture(str(png_file), icon_left, icon_top, icon_width, icon_height)
                        print(f"버튼 아이콘 추가 성공: {icon_class}")
                    else:
                        print(f"버튼 PNG 변환 실패: {icon_class}")
                else:
                    print(f"버튼 SVG 다운로드 실패: {icon_class}")
            
            return button
            
        except Exception as e:
            print(f"버튼 생성 오류: {e}")
            return None
    
    def create_feature_card(self, slide, title, description, icon_class, x, y, width, height):
        """기능 카드 생성 (이미지 아이콘 버전)"""
        try:
            # 카드 배경 생성
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            
            # 배경 색상 설정 (gray-50)
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(249, 250, 251)
            
            # 테두리 설정
            line = card.line
            line.color.rgb = RGBColor(229, 231, 235)
            
            # 텍스트 추가
            text_frame = card.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)
            text_frame.margin_top = Inches(0.2)
            text_frame.margin_bottom = Inches(0.2)
            
            # 제목
            p1 = text_frame.paragraphs[0]
            p1.text = title
            p1.alignment = PP_ALIGN.LEFT
            font1 = p1.font
            font1.size = Pt(14)
            font1.bold = True
            font1.color.rgb = RGBColor(31, 41, 55)
            
            # 설명
            p2 = text_frame.add_paragraph()
            p2.text = description
            p2.alignment = PP_ALIGN.LEFT
            font2 = p2.font
            font2.size = Pt(12)
            font2.color.rgb = RGBColor(107, 114, 128)
            
            # 아이콘 이미지 추가
            if icon_class:
                svg_file = self.download_fontawesome_svg(icon_class, '#1f2937')
                if svg_file:
                    png_file = self.svg_to_png_with_html2image(svg_file, 20)
                    if png_file and png_file.exists():
                        # 아이콘 이미지 추가
                        icon_left = Inches(x + 0.1)
                        icon_top = Inches(y + 0.1)
                        icon_width = Inches(0.2)
                        icon_height = Inches(0.2)
                        
                        slide.shapes.add_picture(str(png_file), icon_left, icon_top, icon_width, icon_height)
                        print(f"기능 카드 아이콘 추가 성공: {icon_class}")
                    else:
                        print(f"기능 카드 PNG 변환 실패: {icon_class}")
                else:
                    print(f"기능 카드 SVG 다운로드 실패: {icon_class}")
            
            return card
            
        except Exception as e:
            print(f"기능 카드 생성 오류: {e}")
            return None
    
    def create_icon_circle(self, slide, icon_class, x, y, size):
        """아이콘 원형 배경 생성 (이미지 버전)"""
        try:
            # 원형 배경 생성
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x), Inches(y), Inches(size), Inches(size)
            )
            
            # 배경 색상 설정 (blue-100)
            fill = circle.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(219, 234, 254)
            
            # 아이콘 다운로드 및 변환
            if icon_class:
                svg_file = self.download_fontawesome_svg(icon_class, '#2563eb')
                if svg_file:
                    png_file = self.svg_to_png_with_html2image(svg_file, 48)
                    if png_file and png_file.exists():
                        # 아이콘 이미지 추가
                        icon_left = Inches(x + size * 0.2)
                        icon_top = Inches(y + size * 0.2)
                        icon_width = Inches(size * 0.6)
                        icon_height = Inches(size * 0.6)
                        
                        slide.shapes.add_picture(str(png_file), icon_left, icon_top, icon_width, icon_height)
                        print(f"아이콘 이미지 추가 성공: {icon_class}")
                    else:
                        print(f"PNG 변환 실패: {icon_class}")
                        # 폴백: 텍스트 아이콘
                        self.add_text_icon_to_circle(circle, icon_class)
                else:
                    print(f"SVG 다운로드 실패: {icon_class}")
                    # 폴백: 텍스트 아이콘
                    self.add_text_icon_to_circle(circle, icon_class)
            else:
                # 아이콘 클래스가 없는 경우 기본 텍스트
                self.add_text_icon_to_circle(circle, None)
            
            return circle
            
        except Exception as e:
            print(f"아이콘 원형 생성 오류: {e}")
            return None
    
    def add_text_icon_to_circle(self, circle, icon_class):
        """원형에 텍스트 아이콘 추가 (폴백)"""
        try:
            if icon_class:
                icon_text = self.get_icon_text(icon_class)
            else:
                icon_text = '●'
            
            text_frame = circle.text_frame
            text_frame.text = icon_text
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # 텍스트 스타일 설정
            p = text_frame.paragraphs[0]
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(37, 99, 235)  # 파란색
            p.font.bold = True
            
        except Exception as e:
            print(f"텍스트 아이콘 추가 실패: {e}")
    
    def get_icon_text(self, icon_class):
        """FontAwesome 아이콘을 텍스트로 변환 (폴백용)"""
        clean_class = icon_class.replace('fas ', '').replace('fab ', '').replace('far ', '')
        icon_map = {
            'fa-react': '⚛',
            'fa-js': 'JS',
            'fa-css3': 'CSS',
            'fa-database': '🗄',
            'fa-server': '🖥',
            'fa-github': '🐙',
            'fa-globe': '🌐',
            'fa-history': '📚',
            'fa-bullseye': '🎯',
            'fa-star': '⭐',
            'fa-users': '👥',
            'fa-graduation-cap': '🎓',
            'fa-project-diagram': '📊',
            'fa-mobile-alt': '📱'
        }
        return icon_map.get(clean_class, '●')
    
    def parse_01_html(self, soup, slide):
        """01.html (메인 페이지) 파싱"""
        # 제목 추가 (text-5xl = 48px)
        title_element = soup.find('h1', class_='title')
        if title_element:
            self.create_text_box(
                slide, title_element.get_text(),
                1, 0.8, 11, 1,
                {'font-size': '48px', 'color': '#2563eb', 'font-weight': '900', 'text-align': 'center'}
            )
        
        # 부제목 추가 (text-3xl = 30px)
        subtitle_element = soup.find('h2', class_='subtitle')
        if subtitle_element:
            self.create_text_box(
                slide, subtitle_element.get_text(),
                1, 1.8, 11, 0.6,
                {'font-size': '30px', 'color': '#1e40af', 'font-weight': '700', 'text-align': 'center'}
            )
        
        # 개발 기간 섹션
        period_label = soup.find('p', string=lambda text: text and '개발 기간' in text)
        if period_label:
            period_value = period_label.find_next_sibling('p')
            if period_value:
                # 라벨
                self.create_text_box(
                    slide, "개발 기간",
                    1, 2.6, 11, 0.3,
                    {'font-size': '20px', 'color': '#6b7280', 'text-align': 'center'}
                )
                # 값
                self.create_text_box(
                    slide, period_value.get_text(),
                    1, 2.9, 11, 0.5,
                    {'font-size': '24px', 'color': '#000000', 'font-weight': '500', 'text-align': 'center'}
                )
        
        # 기술 스택 섹션
        tech_label = soup.find('p', string=lambda text: text and '주요 기술 스택' in text)
        if tech_label:
            # 기술 스택 라벨
            self.create_text_box(
                slide, "주요 기술 스택",
                1, 3.6, 11, 0.3,
                {'font-size': '20px', 'color': '#6b7280', 'text-align': 'center'}
            )
            
            # 기술 스택 컨테이너 찾기
            tech_container = tech_label.find_next('div', class_='flex')
            if tech_container:
                tech_items = tech_container.find_all('div', class_='tech-stack')
                
                # 기술 스택 박스들 (5개를 2행으로 배치)
                start_x = 2.5
                start_y = 4.1
                for i, tech_item in enumerate(tech_items[:5]):
                    tech_text = tech_item.get_text().strip()
                    # 아이콘 추출 및 캡처
                    icon_elem = tech_item.find('i')
                    icon_path = None
                    icon_class = None
                    if icon_elem:
                        icon_classes = icon_elem.get('class', [])
                        for cls in icon_classes:
                            if cls.startswith('fa-'):
                                icon_class = cls
                                break
                    
                    row = i // 3
                    col = i % 3
                    x = start_x + col * 2.8
                    y = start_y + row * 0.8
                    
                    self.create_tech_stack_box(slide, tech_text, icon_class, x, y)
        
        # 링크 버튼들
        links_section = soup.find('div', class_='flex justify-center space-x-8')
        if links_section:
            buttons = links_section.find_all('a', class_='link-button')
            
            # GitHub 버튼 (첫 번째)
            if len(buttons) > 0:
                github_button = buttons[0]
                github_text = github_button.find('span').get_text() if github_button.find('span') else 'GitHub'
                github_icon = github_button.find('i', class_='fab fa-github')
                self.create_button(
                    slide, github_text, 'fab fa-github',
                    4.5, 6.2, 2.5, 0.6,
                    RGBColor(31, 41, 55), RGBColor(255, 255, 255)  # gray-800, white
                )
            
            # 배포 사이트 버튼 (두 번째)
            if len(buttons) > 1:
                deploy_button = buttons[1]
                deploy_text = deploy_button.find('span').get_text() if deploy_button.find('span') else '배포 사이트'
                deploy_icon = deploy_button.find('i', class_='fas fa-globe')
                self.create_button(
                    slide, deploy_text, 'fas fa-globe',
                    7.5, 6.2, 2.5, 0.6,
                    RGBColor(37, 99, 235), RGBColor(255, 255, 255)  # blue-600, white
                )
        
        # 날짜 추가 (우측 하단)
        date_element = soup.find('div', class_='absolute bottom-8 right-8')
        if date_element:
            date_text = date_element.find('p').get_text() if date_element.find('p') else '2025.09.11'
            self.create_text_box(
                slide, date_text,
                10.5, 5.5, 2, 0.3,
                {'font-size': '14px', 'color': '#9ca3af', 'text-align': 'right'}
            )
    
    def parse_02_html(self, soup, slide):
        """02.html (프로젝트 개요) 파싱"""
        # 제목 추가
        title_element = soup.find('h1', class_='section-title')
        if title_element:
            self.create_text_box(
                slide, title_element.get_text(),
                1, 0.3, 11, 0.8,
                {'font-size': '36px', 'color': '#2563eb', 'font-weight': '800', 'text-align': 'left'}
            )
        
        # 배경 섹션
        background_section = soup.find('h2', string=lambda text: text and '배경' in text)
        if background_section:
            # 아이콘 찾기 및 캡처
            parent_div = background_section.find_parent('div', class_='flex')
            icon_path = None
            if parent_div:
                icon_elem = parent_div.find('i')
                icon_class = None
                if icon_elem:
                    icon_classes = icon_elem.get('class', [])
                    for cls in icon_classes:
                        if cls.startswith('fa-'):
                            icon_class = cls
                            break
            
            # 아이콘 원형
            self.create_icon_circle(slide, icon_class, 0.5, 1.2, 0.6)
            
            # 제목
            self.create_text_box(
                slide, "배경",
                1.3, 1.2, 10, 0.5,
                {'font-size': '24px', 'color': '#1f2937', 'font-weight': 'bold', 'text-align': 'left'}
            )
            
            # 내용
            background_content = background_section.find_next('p')
            if background_content:
                self.create_text_box(
                    slide, background_content.get_text(),
                    1.3, 1.7, 10, 1,
                    {'font-size': '16px', 'color': '#6b7280', 'text-align': 'left'}
                )
        
        # 목적 섹션
        purpose_section = soup.find('h2', string=lambda text: text and '목적' in text)
        if purpose_section:
            # 아이콘 찾기 및 캡처
            parent_div = purpose_section.find_parent('div', class_='flex')
            icon_path = None
            if parent_div:
                icon_elem = parent_div.find('i')
                icon_class = None
                if icon_elem:
                    icon_classes = icon_elem.get('class', [])
                    for cls in icon_classes:
                        if cls.startswith('fa-'):
                            icon_class = cls
                            break
            
            # 아이콘 원형
            self.create_icon_circle(slide, icon_class, 0.5, 2.9, 0.6)
            
            # 제목
            self.create_text_box(
                slide, "목적",
                1.3, 2.9, 10, 0.5,
                {'font-size': '24px', 'color': '#1f2937', 'font-weight': 'bold', 'text-align': 'left'}
            )
            
            # 내용
            purpose_content = purpose_section.find_next('p')
            if purpose_content:
                self.create_text_box(
                    slide, purpose_content.get_text(),
                    1.3, 3.4, 10, 1,
                    {'font-size': '16px', 'color': '#6b7280', 'text-align': 'left'}
                )
        
        # 주요 특징 섹션
        features_section = soup.find('h2', string=lambda text: text and '주요 특징' in text)
        if features_section:
            # 아이콘 찾기 및 캡처
            parent_div = features_section.find_parent('div', class_='flex')
            icon_path = None
            if parent_div:
                icon_elem = parent_div.find('i')
                icon_class = None
                if icon_elem:
                    icon_classes = icon_elem.get('class', [])
                    for cls in icon_classes:
                        if cls.startswith('fa-'):
                            icon_class = cls
                            break
            
            # 아이콘 원형
            self.create_icon_circle(slide, icon_class, 0.5, 4.6, 0.6)
            
            # 제목
            self.create_text_box(
                slide, "주요 특징",
                1.3, 4.6, 10, 0.5,
                {'font-size': '24px', 'color': '#1f2937', 'font-weight': 'bold', 'text-align': 'left'}
            )
            
            # 특징 카드들
            feature_cards = soup.find_all('div', class_='feature-card')
            if feature_cards:
                # 2x2 그리드로 배치
                start_x = 1.3
                start_y = 5.2
                card_width = 4.5
                card_height = 1.0
                
                for i, card in enumerate(feature_cards[:4]):
                    row = i // 2
                    col = i % 2
                    x = start_x + col * 5.0
                    y = start_y + row * 1.2
                    
                    title_elem = card.find('h3')
                    desc_elem = card.find('p')
                    icon_elem = card.find('i')
                    
                    if title_elem and desc_elem:
                        title = title_elem.get_text()
                        description = desc_elem.get_text()
                        
                        # 아이콘 캡처
                        icon_path = None
                        if icon_elem:
                            icon_classes = icon_elem.get('class', [])
                            for cls in icon_classes:
                                if cls.startswith('fa-'):
                                    icon_class = cls
                                    break
                        
                        self.create_feature_card(slide, title, description, icon_class, x, y, card_width, card_height)
        
        # 푸터
        footer_element = soup.find('div', class_='absolute bottom-8 right-8')
        if footer_element:
            footer_text = footer_element.find('p').get_text() if footer_element.find('p') else '개발 프로젝트: 디지털 창작소 웹사이트'
            self.create_text_box(
                slide, footer_text,
                8, 6.8, 4, 0.3,
                {'font-size': '12px', 'color': '#9ca3af', 'text-align': 'right'}
            )
    
    def parse_html_to_pptx(self):
        """HTML을 파싱하여 PPTX로 변환"""
        try:
            # HTML 파일 읽기
            with open(self.html_file, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # 새 프레젠테이션 생성
            prs = Presentation()
            
            # 슬라이드 크기 설정 (16:9 비율)
            prs.slide_width = Inches(13.33)  # 1920px
            prs.slide_height = Inches(7.5)   # 1080px
            
            # 빈 슬라이드 추가
            slide_layout = prs.slide_layouts[6]  # 빈 레이아웃
            slide = prs.slides.add_slide(slide_layout)
            
            # HTML 파일 타입에 따라 다른 파싱 로직 적용
            if '01.html' in str(self.html_file):
                print("01.html (메인 페이지) 파싱 중...")
                self.parse_01_html(soup, slide)
            elif '02.html' in str(self.html_file):
                print("02.html (프로젝트 개요) 파싱 중...")
                self.parse_02_html(soup, slide)
            else:
                print("알 수 없는 HTML 파일 형식입니다.")
                return False
            
            # PPTX 파일 저장
            prs.save(self.output_path)
            print(f"PPTX 파일 저장 완료: {self.output_path}")
            
            return True
            
        except Exception as e:
            print(f"HTML 파싱 오류: {e}")
            import traceback
            traceback.print_exc()
            return False
    
    def convert(self):
        """전체 변환 프로세스 실행"""
        try:
            print(f"HTML 파일 변환 시작: {self.html_file}")
            
            # 임시 디렉토리 설정
            self.setup_temp_directory()
            
            # HTML을 PPTX로 변환
            success = self.parse_html_to_pptx()
            
            return success
            
        except Exception as e:
            print(f"변환 프로세스 오류: {e}")
            return False
        
        finally:
            # 임시 디렉토리 정리
            self.cleanup_temp_directory()

def convert_folder_to_pptx(html_folder, output_path):
    """폴더 내 모든 HTML 파일을 하나의 PPTX로 변환"""
    try:
        html_folder = Path(html_folder)
        if not html_folder.exists():
            print(f"HTML 폴더가 존재하지 않습니다: {html_folder}")
            return False
        
        # HTML 파일들 찾기
        html_files = list(html_folder.glob("*.html"))
        if not html_files:
            print(f"HTML 파일이 없습니다: {html_folder}")
            return False
        
        print(f"발견된 HTML 파일 {len(html_files)}개:")
        for html_file in html_files:
            print(f"  - {html_file.name}")
        
        # PPTX 프레젠테이션 생성
        prs = Presentation()
        
        # 각 HTML 파일을 슬라이드로 변환
        for i, html_file in enumerate(html_files):
            print(f"\n--- {html_file.name} 변환 중 ({i+1}/{len(html_files)}) ---")
            
            try:
                # 슬라이드 추가
                slide_layout = prs.slide_layouts[6]  # 빈 슬라이드
                slide = prs.slides.add_slide(slide_layout)
                
                # HTML 파일 읽기
                with open(html_file, 'r', encoding='utf-8') as f:
                    html_content = f.read()
                
                soup = BeautifulSoup(html_content, 'html.parser')
                
                # HTML 파일 타입에 따라 파싱
                converter = HTMLEditablePPTXConverterV6(str(html_file), "")
                converter.setup_temp_directory()  # 임시 디렉토리 설정
                
                if '01.html' in html_file.name:
                    converter.parse_01_html(soup, slide)
                elif '02.html' in html_file.name:
                    converter.parse_02_html(soup, slide)
                else:
                    # 기본 파싱 (제목과 내용만)
                    parse_generic_html(converter, soup, slide)
                
                print(f"✅ {html_file.name} 변환 완료")
                
            except Exception as e:
                print(f"❌ {html_file.name} 변환 실패: {e}")
                continue
        
        # PPTX 파일 저장
        prs.save(output_path)
        print(f"\n✅ 모든 HTML 파일이 하나의 PPTX로 변환 완료!")
        print(f"출력 파일: {output_path}")
        return True
        
    except Exception as e:
        print(f"폴더 변환 오류: {e}")
        return False

def parse_generic_html(converter, soup, slide):
    """일반적인 HTML 파일 파싱"""
    try:
        # 제목 찾기
        title = soup.find('h1') or soup.find('title')
        if title:
            title_text = title.get_text().strip()
            converter.create_text_box(
                slide, title_text,
                1.0, 1.0, 8.0, 1.0,
                {'font-size': '32px', 'color': '#1f2937', 'font-weight': 'bold', 'text-align': 'center'}
            )
        
        # 본문 내용 찾기
        content = soup.find('main') or soup.find('body')
        if content:
            paragraphs = content.find_all('p')
            y_pos = 2.5
            for p in paragraphs[:5]:  # 최대 5개 문단
                text = p.get_text().strip()
                if text:
                    converter.create_text_box(
                        slide, text,
                        1.0, y_pos, 8.0, 0.8,
                        {'font-size': '16px', 'color': '#374151', 'text-align': 'left'}
                    )
                    y_pos += 1.0
        
    except Exception as e:
        print(f"일반 HTML 파싱 오류: {e}")

def main():
    # 설정
    html_folder = r"C:\Project\gigabitamin\genspark\dcs_site\html"
    output_path = r"C:\Project\gigabitamin\genspark\dcs_site\html\all_pages_editable_v6.pptx"
    
    print("HTML 폴더 to Editable PPTX 변환기 V6 시작")
    print(f"HTML 폴더: {html_folder}")
    print(f"출력 파일: {output_path}")
    print("-" * 50)
    
    # 폴더 변환 실행
    success = convert_folder_to_pptx(html_folder, output_path)
    
    if success:
        print("-" * 50)
        print("변환 완료!")
        print(f"출력 파일: {output_path}")
        print("변환된 PPTX 파일에서 텍스트, 도형, 이미지를 개별적으로 편집할 수 있습니다.")
    else:
        print("-" * 50)
        print("변환 실패!")
        sys.exit(1)

if __name__ == "__main__":
    main()