#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
HTML to PPTX Converter
HTML 파일들을 스타일을 유지한 채로 PPTX 파일로 변환하는 스크립트
"""

import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from html2image import Html2Image
from PIL import Image
import tempfile
import shutil

class HTMLToPPTXConverter:
    def __init__(self, html_dir, output_path):
        self.html_dir = Path(html_dir)
        self.output_path = Path(output_path)
        self.hti = Html2Image()
        self.temp_dir = None
        
    def setup_temp_directory(self):
        """임시 디렉토리 설정"""
        self.temp_dir = Path(tempfile.mkdtemp())
        print(f"임시 디렉토리 생성: {self.temp_dir}")
        
    def cleanup_temp_directory(self):
        """임시 디렉토리 정리"""
        if self.temp_dir and self.temp_dir.exists():
            shutil.rmtree(self.temp_dir)
            print("임시 디렉토리 정리 완료")
    
    def convert_html_to_image(self, html_file, slide_number):
        """HTML 파일을 이미지로 변환"""
        try:
            # HTML 파일 경로
            html_path = self.html_dir / html_file
            
            # 출력 이미지 경로
            output_image = self.temp_dir / f"slide_{slide_number:02d}.png"
            
            print(f"변환 중: {html_file} -> {output_image.name}")
            
            # HTML을 이미지로 변환 (더 큰 크기로 설정하여 전체 페이지 캡처)
            self.hti.screenshot(
                html_file=str(html_path),
                save_as=f"slide_{slide_number:02d}.png",
                size=(1280, 1440)  # 높이를 더 크게 설정
            )
            
            # 생성된 이미지 파일을 임시 디렉토리로 이동
            generated_image = Path(f"slide_{slide_number:02d}.png")
            if generated_image.exists():
                shutil.move(str(generated_image), str(output_image))
                
                # 이미지 크기 조정 (720 높이로 리사이즈)
                self.resize_image_to_fit(output_image)
            
            return output_image
            
        except Exception as e:
            print(f"HTML 변환 오류 ({html_file}): {e}")
            return None
    
    def resize_image_to_fit(self, image_path):
        """이미지를 1280x720에 맞게 리사이즈"""
        try:
            from PIL import Image
            
            with Image.open(image_path) as img:
                # 원본 이미지 크기
                original_width, original_height = img.size
                
                print(f"원본 이미지 크기: {original_width}x{original_height}")
                
                # 목표 크기
                target_width = 1280
                target_height = 720
                
                # 원본 비율 계산
                original_ratio = original_width / original_height
                target_ratio = target_width / target_height
                
                if original_ratio > target_ratio:
                    # 원본이 더 넓은 경우 - 높이를 기준으로 리사이즈
                    new_height = target_height
                    new_width = int(target_height * original_ratio)
                else:
                    # 원본이 더 높은 경우 - 너비를 기준으로 리사이즈
                    new_width = target_width
                    new_height = int(target_width / original_ratio)
                
                print(f"리사이즈 크기: {new_width}x{new_height}")
                
                # 리사이즈
                resized_img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
                
                # 1280x720으로 크롭 (상단부터 시작)
                if new_width > target_width:
                    # 너비가 더 긴 경우 - 중앙에서 크롭
                    left = (new_width - target_width) // 2
                    right = left + target_width
                else:
                    left = 0
                    right = target_width
                
                if new_height > target_height:
                    # 높이가 더 긴 경우 - 상단부터 크롭 (위쪽이 잘리지 않도록)
                    top = 0
                    bottom = target_height
                else:
                    top = 0
                    bottom = new_height
                
                print(f"크롭 영역: left={left}, top={top}, right={right}, bottom={bottom}")
                
                cropped_img = resized_img.crop((left, top, right, bottom))
                
                # 최종 크기가 1280x720이 되도록 패딩 추가 (필요시)
                if cropped_img.size != (target_width, target_height):
                    # 흰색 배경으로 새 이미지 생성
                    final_img = Image.new('RGB', (target_width, target_height), 'white')
                    # 중앙에 이미지 붙이기
                    paste_x = (target_width - cropped_img.width) // 2
                    paste_y = (target_height - cropped_img.height) // 2
                    final_img.paste(cropped_img, (paste_x, paste_y))
                    cropped_img = final_img
                
                # 저장
                cropped_img.save(image_path, 'PNG', quality=95)
                
                print(f"이미지 리사이즈 완료: {image_path.name} -> {cropped_img.size}")
                
        except Exception as e:
            print(f"이미지 리사이즈 오류: {e}")
    
    def create_pptx(self, image_files):
        """이미지 파일들로부터 PPTX 생성"""
        try:
            # 새 프레젠테이션 생성
            prs = Presentation()
            
            # 슬라이드 크기 설정 (16:9 비율)
            prs.slide_width = Inches(13.33)  # 1280px
            prs.slide_height = Inches(7.5)   # 720px
            
            for i, image_file in enumerate(image_files):
                if image_file and image_file.exists():
                    # 빈 슬라이드 추가
                    slide_layout = prs.slide_layouts[6]  # 빈 레이아웃
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # 이미지 추가
                    left = Inches(0)
                    top = Inches(0)
                    width = Inches(13.33)
                    height = Inches(7.5)
                    
                    slide.shapes.add_picture(str(image_file), left, top, width, height)
                    
                    print(f"슬라이드 {i+1} 추가 완료")
                else:
                    print(f"이미지 파일 없음: {image_file}")
            
            # PPTX 파일 저장
            prs.save(self.output_path)
            print(f"PPTX 파일 저장 완료: {self.output_path}")
            
        except Exception as e:
            print(f"PPTX 생성 오류: {e}")
            raise
    
    def convert(self):
        """전체 변환 프로세스 실행"""
        try:
            # HTML 파일 목록 가져오기
            html_files = sorted([f for f in self.html_dir.glob("*.html")])
            
            if not html_files:
                print("HTML 파일을 찾을 수 없습니다.")
                return False
            
            print(f"발견된 HTML 파일: {len(html_files)}개")
            for html_file in html_files:
                print(f"  - {html_file.name}")
            
            # 임시 디렉토리 설정
            self.setup_temp_directory()
            
            # HTML 파일들을 이미지로 변환
            image_files = []
            for i, html_file in enumerate(html_files, 1):
                image_file = self.convert_html_to_image(html_file.name, i)
                image_files.append(image_file)
            
            # PPTX 생성
            self.create_pptx(image_files)
            
            return True
            
        except Exception as e:
            print(f"변환 프로세스 오류: {e}")
            return False
        
        finally:
            # 임시 디렉토리 정리
            self.cleanup_temp_directory()

def main():
    # 설정
    html_dir = r"C:\Project\gigabitamin\notion\doc_ppt\hearth_chat"
    output_path = r"C:\Project\gigabitamin\notion\doc_ppt\hearth_chat\hearth_chat_presentation.pptx"
    
    print("HTML to PPTX 변환기 시작")
    print(f"HTML 디렉토리: {html_dir}")
    print(f"출력 파일: {output_path}")
    print("-" * 50)
    
    # 변환기 생성 및 실행
    converter = HTMLToPPTXConverter(html_dir, output_path)
    success = converter.convert()
    
    if success:
        print("-" * 50)
        print("변환 완료!")
        print(f"출력 파일: {output_path}")
    else:
        print("-" * 50)
        print("변환 실패!")
        sys.exit(1)

if __name__ == "__main__":
    main()
