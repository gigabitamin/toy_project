#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Universal HTML to Editable PPTX Converter
모든 HTML 파일을 일관되게 변환하는 범용 변환기
"""

import os
import sys
import requests
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup
import re
import tempfile
import shutil
from html2image import Html2Image
from PIL import Image

class UniversalHTMLToPPTXConverter:
    def __init__(self, html_file, output_path):
        self.html_file = html_file
        self.output_path = output_path
        self.temp_dir = None
        self.hti = Html2Image()
        
    def setup_temp_directory(self):
        """임시 디렉토리 생성"""
        self.temp_dir = Path(tempfile.mkdtemp())
        
    def cleanup_temp_directory(self):
        """임시 디렉토리 정리"""
        if self.temp_dir and self.temp_dir.exists():
            shutil.rmtree(self.temp_dir)
    
    def download_fontawesome_svg(self, icon_class, color='#2563eb'):
        """FontAwesome 아이콘 SVG 다운로드"""
        try:
            clean_class = icon_class.replace('fas ', '').replace('fab ', '').replace('far ', '').replace('fa-', '')
            
            # 다양한 경로 시도
            paths = [
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/6.x/svgs/solid/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/6.x/svgs/brands/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/6.x/svgs/regular/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/5.x/svgs/solid/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/5.x/svgs/brands/{clean_class}.svg",
                f"https://raw.githubusercontent.com/FortAwesome/Font-Awesome/5.x/svgs/regular/{clean_class}.svg",
            ]
            
            for path in paths:
                try:
                    response = requests.get(path, timeout=10)
                    if response.status_code == 200:
                        svg_content = response.text
                        
                        # 색상 적용
                        svg_content = svg_content.replace('fill="currentColor"', f'fill="{color}"')
                        svg_content = svg_content.replace('fill="#000"', f'fill="{color}"')
                        svg_content = svg_content.replace('fill="black"', f'fill="{color}"')
                        
                        # SVG 파일 저장
                        svg_file = self.temp_dir / f"{clean_class}.svg"
                        with open(svg_file, 'w', encoding='utf-8') as f:
                            f.write(svg_content)
                        
                        return svg_file
                        
                except Exception as e:
                    continue
            
            return None
            
        except Exception as e:
            return None
    
    def svg_to_png_with_html2image(self, svg_file, size=64):
        """SVG를 HTML2Image로 PNG 변환"""
        try:
            with open(svg_file, 'r', encoding='utf-8') as f:
                svg_content = f.read()
            
            # HTML 생성
            html_content = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <style>
                    body {{
                        margin: 0;
                        padding: 0;
                        width: {size}px;
                        height: {size}px;
                        display: flex;
                        align-items: center;
                        justify-content: center;
                        background: transparent;
                    }}
                    svg {{
                        width: {size}px;
                        height: {size}px;
                    }}
                </style>
            </head>
            <body>
                {svg_content}
            </body>
            </html>
            """
            
            # 임시 HTML 파일 생성
            temp_html = self.temp_dir / f"{svg_file.stem}_temp.html"
            with open(temp_html, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            # PNG 파일 경로
            png_file = self.temp_dir / f"{svg_file.stem}.png"
            
            # HTML2Image로 변환
            self.hti.screenshot(
                html_file=str(temp_html),
                save_as=f"{svg_file.stem}.png",
                size=(size, size)
            )
            
            # 생성된 PNG 파일을 임시 디렉토리로 이동
            generated_png = Path(f"{svg_file.stem}.png")
            if generated_png.exists():
                shutil.move(str(generated_png), str(png_file))
                return png_file
            else:
                return None
                
        except Exception as e:
            return None
    
    def create_text_box(self, slide, text, x, y, width, height, styles=None):
        """텍스트 박스 생성"""
        try:
            # 텍스트 박스 추가
            textbox = slide.shapes.add_textbox(
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            
            # 텍스트 프레임 설정
            text_frame = textbox.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            # 텍스트 추가
            p = text_frame.paragraphs[0]
            p.text = text
            
            # 스타일 적용
            if styles:
                if 'font-size' in styles:
                    font_size = self.parse_font_size(styles['font-size'])
                    p.font.size = Pt(font_size)
                
                if 'color' in styles:
                    color = self.parse_css_color(styles['color'])
                    p.font.color.rgb = color
                
                if 'font-weight' in styles:
                    if 'bold' in styles['font-weight']:
                        p.font.bold = True
                
                if 'text-align' in styles:
                    alignment = self.get_text_alignment(styles['text-align'])
                    p.alignment = alignment
            
            return textbox
            
        except Exception as e:
            print(f"텍스트 박스 생성 오류: {e}")
            return None
    
    def create_icon_image(self, slide, icon_class, x, y, size=0.5, color='#2563eb'):
        """아이콘 이미지 생성"""
        try:
            if not icon_class:
                return None
                
            svg_file = self.download_fontawesome_svg(icon_class, color)
            if svg_file:
                png_file = self.svg_to_png_with_html2image(svg_file, 48)
                if png_file and png_file.exists():
                    # 아이콘 이미지 추가
                    icon_left = Inches(x)
                    icon_top = Inches(y)
                    icon_width = Inches(size)
                    icon_height = Inches(size)
                    
                    slide.shapes.add_picture(str(png_file), icon_left, icon_top, icon_width, icon_height)
                    return True
            
            return False
            
        except Exception as e:
            print(f"아이콘 이미지 생성 오류: {e}")
            return False
    
    def create_card(self, slide, title, content, x, y, width, height, icon_class=None):
        """카드 생성"""
        try:
            # 카드 배경 생성
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            
            # 배경 색상 설정
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(249, 250, 251)
            
            # 테두리 설정
            line = card.line
            line.color.rgb = RGBColor(229, 231, 235)
            
            # 텍스트 추가
            text_frame = card.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)
            text_frame.margin_top = Inches(0.2)
            text_frame.margin_bottom = Inches(0.2)
            
            # 제목
            p1 = text_frame.paragraphs[0]
            p1.text = title
            p1.alignment = PP_ALIGN.LEFT
            font1 = p1.font
            font1.size = Pt(14)
            font1.bold = True
            font1.color.rgb = RGBColor(31, 41, 55)
            
            # 내용
            if content:
                p2 = text_frame.add_paragraph()
                p2.text = content
                p2.alignment = PP_ALIGN.LEFT
                font2 = p2.font
                font2.size = Pt(12)
                font2.color.rgb = RGBColor(107, 114, 128)
            
            # 아이콘 추가
            if icon_class:
                self.create_icon_image(slide, icon_class, x + 0.1, y + 0.1, 0.2, '#1f2937')
            
            return card
            
        except Exception as e:
            print(f"카드 생성 오류: {e}")
            return None
    
    def parse_css_color(self, color_str):
        """CSS 색상을 RGB로 변환"""
        if color_str.startswith('#'):
            hex_color = color_str[1:]
            if len(hex_color) == 3:
                hex_color = ''.join([c*2 for c in hex_color])
            return RGBColor(
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16),
                int(hex_color[4:6], 16)
            )
        elif color_str.startswith('rgb'):
            rgb_values = re.findall(r'\d+', color_str)
            if len(rgb_values) >= 3:
                return RGBColor(
                    int(rgb_values[0]),
                    int(rgb_values[1]),
                    int(rgb_values[2])
                )
        return RGBColor(0, 0, 0)
    
    def parse_font_size(self, font_size_str):
        """폰트 크기 파싱"""
        if 'px' in font_size_str:
            return int(font_size_str.replace('px', ''))
        elif 'pt' in font_size_str:
            return int(font_size_str.replace('pt', ''))
        return 12
    
    def get_text_alignment(self, align_str):
        """텍스트 정렬 파싱"""
        align_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        return align_map.get(align_str, PP_ALIGN.LEFT)
    
    def extract_icon_class(self, element):
        """요소에서 아이콘 클래스 추출"""
        if not element:
            return None
            
        # i 태그에서 클래스 찾기
        icon_elem = element.find('i')
        if icon_elem:
            classes = icon_elem.get('class', [])
            for cls in classes:
                if cls.startswith('fa-'):
                    return cls
        
        # 요소 자체에서 클래스 찾기
        classes = element.get('class', [])
        for cls in classes:
            if cls.startswith('fa-'):
                return cls
        
        return None
    
    def parse_html_universal(self, soup, slide):
        """범용 HTML 파싱"""
        try:
            y_pos = 0.5
            
            # 제목 찾기 (h1, h2, title 순서)
            title_element = soup.find('h1') or soup.find('h2') or soup.find('title')
            if title_element:
                title_text = title_element.get_text().strip()
                if title_text:
                    self.create_text_box(
                        slide, title_text,
                        0.5, y_pos, 9.0, 1.0,
                        {'font-size': '32px', 'color': '#1f2937', 'font-weight': 'bold', 'text-align': 'center'}
                    )
                    y_pos += 1.2
            
            # 섹션들 찾기
            sections = soup.find_all(['section', 'div'], class_=re.compile(r'section|container|content'))
            
            for section in sections[:6]:  # 최대 6개 섹션
                # 섹션 제목 찾기
                section_title = section.find(['h2', 'h3', 'h4'])
                if section_title:
                    title_text = section_title.get_text().strip()
                    if title_text:
                        # 아이콘 찾기
                        icon_class = self.extract_icon_class(section_title)
                        
                        # 아이콘 이미지 추가
                        if icon_class:
                            self.create_icon_image(slide, icon_class, 0.5, y_pos, 0.4, '#2563eb')
                            title_x = 1.0
                        else:
                            title_x = 0.5
                        
                        # 섹션 제목
                        self.create_text_box(
                            slide, title_text,
                            title_x, y_pos, 8.0, 0.6,
                            {'font-size': '24px', 'color': '#1f2937', 'font-weight': 'bold', 'text-align': 'left'}
                        )
                        y_pos += 0.8
                
                # 섹션 내용 찾기
                content_elements = section.find_all(['p', 'div'], class_=re.compile(r'text|content|description'))
                
                for content_elem in content_elements[:3]:  # 최대 3개 내용
                    content_text = content_elem.get_text().strip()
                    if content_text and len(content_text) > 10:  # 의미있는 내용만
                        self.create_text_box(
                            slide, content_text,
                            0.5, y_pos, 9.0, 0.8,
                            {'font-size': '16px', 'color': '#374151', 'text-align': 'left'}
                        )
                        y_pos += 1.0
                
                # 카드들 찾기
                cards = section.find_all(['div'], class_=re.compile(r'card|box|item'))
                
                for i, card in enumerate(cards[:4]):  # 최대 4개 카드
                    card_title = card.find(['h3', 'h4', 'h5'])
                    card_content = card.find(['p', 'div'], class_=re.compile(r'text|content|description'))
                    
                    if card_title:
                        title_text = card_title.get_text().strip()
                        content_text = card_content.get_text().strip() if card_content else ""
                        
                        # 아이콘 찾기
                        icon_class = self.extract_icon_class(card)
                        
                        # 카드 위치 계산 (2x2 그리드)
                        card_x = 0.5 + (i % 2) * 4.5
                        card_y = y_pos + (i // 2) * 1.5
                        
                        self.create_card(
                            slide, title_text, content_text,
                            card_x, card_y, 4.0, 1.2, icon_class
                        )
                
                if cards:
                    y_pos += 2.5  # 카드가 있으면 추가 공간
                else:
                    y_pos += 0.5  # 카드가 없으면 작은 공간
                
                # 슬라이드 높이 제한
                if y_pos > 6.0:
                    break
            
            # 리스트 항목들 찾기
            lists = soup.find_all(['ul', 'ol'])
            for list_elem in lists[:2]:  # 최대 2개 리스트
                items = list_elem.find_all('li')
                for item in items[:5]:  # 최대 5개 항목
                    item_text = item.get_text().strip()
                    if item_text:
                        self.create_text_box(
                            slide, f"• {item_text}",
                            0.5, y_pos, 9.0, 0.4,
                            {'font-size': '14px', 'color': '#374151', 'text-align': 'left'}
                        )
                        y_pos += 0.5
                
                if items:
                    y_pos += 0.3
            
            # 코드 블록 찾기
            code_blocks = soup.find_all(['pre', 'code'])
            for code_block in code_blocks[:2]:  # 최대 2개 코드 블록
                code_text = code_block.get_text().strip()
                if code_text and len(code_text) > 20:
                    # 코드 블록 배경
                    code_bg = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.5), Inches(y_pos), Inches(9.0), Inches(1.0)
                    )
                    code_bg.fill.solid()
                    code_bg.fill.fore_color.rgb = RGBColor(248, 249, 250)
                    
                    # 코드 텍스트
                    self.create_text_box(
                        slide, code_text,
                        0.7, y_pos + 0.1, 8.6, 0.8,
                        {'font-size': '12px', 'color': '#1f2937', 'text-align': 'left'}
                    )
                    y_pos += 1.2
            
        except Exception as e:
            print(f"범용 HTML 파싱 오류: {e}")
    
    def convert(self):
        """HTML을 PPTX로 변환"""
        try:
            # 임시 디렉토리 설정
            self.setup_temp_directory()
            
            # HTML 파일 읽기
            with open(self.html_file, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # PPTX 프레젠테이션 생성
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]  # 빈 슬라이드
            slide = prs.slides.add_slide(slide_layout)
            
            # 범용 파싱 실행
            self.parse_html_universal(soup, slide)
            
            # PPTX 파일 저장
            prs.save(self.output_path)
            print(f"✅ 변환 완료: {Path(self.html_file).name}")
            
            return True
            
        except Exception as e:
            print(f"❌ 변환 실패: {e}")
            return False
        
        finally:
            # 임시 디렉토리 정리
            self.cleanup_temp_directory()

def convert_folder_to_pptx(html_folder, output_path):
    """폴더 내 모든 HTML 파일을 하나의 PPTX로 변환"""
    try:
        html_folder = Path(html_folder)
        if not html_folder.exists():
            print(f"HTML 폴더가 존재하지 않습니다: {html_folder}")
            return False
        
        # HTML 파일들 찾기
        html_files = list(html_folder.glob("*.html"))
        if not html_files:
            print(f"HTML 파일이 없습니다: {html_folder}")
            return False
        
        print(f"발견된 HTML 파일 {len(html_files)}개:")
        for html_file in html_files:
            print(f"  - {html_file.name}")
        
        # PPTX 프레젠테이션 생성
        prs = Presentation()
        
        # 각 HTML 파일을 슬라이드로 변환
        for i, html_file in enumerate(html_files):
            print(f"\n--- {html_file.name} 변환 중 ({i+1}/{len(html_files)}) ---")
            
            try:
                # 슬라이드 추가
                slide_layout = prs.slide_layouts[6]  # 빈 슬라이드
                slide = prs.slides.add_slide(slide_layout)
                
                # 변환기 생성 및 실행
                converter = UniversalHTMLToPPTXConverter(str(html_file), "")
                converter.setup_temp_directory()
                
                # HTML 파일 읽기
                with open(html_file, 'r', encoding='utf-8') as f:
                    html_content = f.read()
                
                soup = BeautifulSoup(html_content, 'html.parser')
                
                # 범용 파싱 실행
                converter.parse_html_universal(soup, slide)
                converter.cleanup_temp_directory()
                
                print(f"✅ {html_file.name} 변환 완료")
                
            except Exception as e:
                print(f"❌ {html_file.name} 변환 실패: {e}")
                continue
        
        # PPTX 파일 저장
        prs.save(output_path)
        print(f"\n✅ 모든 HTML 파일이 하나의 PPTX로 변환 완료!")
        print(f"출력 파일: {output_path}")
        return True
        
    except Exception as e:
        print(f"폴더 변환 오류: {e}")
        return False

def main():
    # 설정
    html_folder = r"C:\Project\gigabitamin\genspark\dcs_site\html"
    output_path = r"C:\Project\gigabitamin\genspark\dcs_site\html\universal_all_pages.pptx"
    
    print("Universal HTML to Editable PPTX 변환기")
    print("=" * 50)
    print(f"HTML 폴더: {html_folder}")
    print(f"출력 파일: {output_path}")
    print("-" * 50)
    
    # 폴더 변환 실행
    success = convert_folder_to_pptx(html_folder, output_path)
    
    if success:
        print("-" * 50)
        print("✅ 변환 완료!")
        print(f"출력 파일: {output_path}")
        print(f"파일 크기: {Path(output_path).stat().st_size:,} bytes")
    else:
        print("❌ 변환 실패!")

if __name__ == "__main__":
    main()

